from django.shortcuts import render, redirect, get_object_or_404
from django.contrib.auth.decorators import login_required
from django.contrib import messages
from django.db.models import Count, Q
from django.db import transaction
from django.core.paginator import Paginator
from django.http import HttpResponse, JsonResponse
from django.views.decorators.http import require_http_methods
from django.urls import reverse
from django.core.exceptions import ValidationError
from apps.infrastructure.models import CustomUser, Company, Employee, Evaluation, RiskResult, Question, Response, Dimension
from apps.presentation.forms import CompanyForm, EmployeeForm, BulkImportForm
from apps.presentation.utils.ecuador_data import CIUDADES_POR_PROVINCIA
from apps.presentation.utils.security import (
    sanitize_string, validate_integer, validate_year, 
    sanitize_sql_like_pattern, validate_company_access, validate_employee_access
)
from apps.application.services.risk_calculator import RiskCalculatorService
from apps.application.services.recommendations_service import RecommendationsService
from datetime import datetime
import openpyxl
from openpyxl.styles import Font, Alignment, PatternFill
from openpyxl.utils import get_column_letter
import csv
import secrets
import string
import json


def user_is_admin(user):
    """Verifica si el usuario es admin o superuser"""
    return user.is_superuser or user.role in [CustomUser.Role.SUPERUSER, CustomUser.Role.COMPANY_ADMIN]


@login_required
def admin_dashboard(request):
    """Dashboard para administradores de empresa"""
    if not user_is_admin(request.user):
        messages.error(request, 'No tienes permisos para acceder a esta página')
        return redirect('employee_dashboard')
    
    # Obtener las empresas que gestiona este admin
    if request.user.is_superuser or request.user.role == CustomUser.Role.SUPERUSER:
        companies = Company.objects.all()
    else:
        companies = request.user.managed_companies.all()
    
    # Estadísticas generales
    total_companies = companies.count()
    total_employees = Employee.objects.filter(company__in=companies).count()
    total_evaluations = Evaluation.objects.filter(
        employee__company__in=companies,
        status=Evaluation.Status.COMPLETED
    ).count()
    
    current_year = datetime.now().year
    evaluations_this_year = Evaluation.objects.filter(
        employee__company__in=companies,
        year=current_year,
        status=Evaluation.Status.COMPLETED
    ).count()
    
    context = {
        'companies': companies,
        'total_companies': total_companies,
        'total_employees': total_employees,
        'total_evaluations': total_evaluations,
        'evaluations_this_year': evaluations_this_year,
        'current_year': current_year,
    }
    
    return render(request, 'admin/dashboard.html', context)


@login_required
def superuser_dashboard(request):
    """Dashboard para superusuarios"""
    if not (request.user.is_superuser or request.user.role == CustomUser.Role.SUPERUSER):
        messages.error(request, 'No tienes permisos para acceder a esta página')
        return redirect('login')
    
    # Estadísticas globales
    total_companies = Company.objects.count()
    total_admins = CustomUser.objects.filter(role=CustomUser.Role.COMPANY_ADMIN).count()
    total_employees = Employee.objects.count()
    total_evaluations = Evaluation.objects.filter(status=Evaluation.Status.COMPLETED).count()
    
    # Obtener una evaluación de ejemplo para el botón (la más reciente completada)
    sample_evaluation = Evaluation.objects.filter(status=Evaluation.Status.COMPLETED).order_by('-date_completed', '-id').first()
    
    current_year = datetime.now().year
    
    context = {
        'total_companies': total_companies,
        'total_admins': total_admins,
        'total_employees': total_employees,
        'total_evaluations': total_evaluations,
        'current_year': current_year,
        'sample_evaluation': sample_evaluation,
    }
    
    return render(request, 'admin/superuser_dashboard.html', context)


@login_required
def company_list(request):
    """Lista de empresas"""
    if not user_is_admin(request.user):
        messages.error(request, 'No tienes permisos para acceder a esta página')
        return redirect('employee_dashboard')
    
    if request.user.is_superuser or request.user.role == CustomUser.Role.SUPERUSER:
        companies = Company.objects.all().annotate(
            employee_count=Count('employees')
        )
    else:
        companies = request.user.managed_companies.all().annotate(
            employee_count=Count('employees')
        )
    
    context = {
        'companies': companies,
    }
    
    return render(request, 'admin/company_list.html', context)


@login_required
def employee_list(request, company_id=None):
    """Lista de empleados con funcionalidades avanzadas"""
    if not user_is_admin(request.user):
        messages.error(request, 'No tienes permisos para acceder a esta página')
        return redirect('employee_dashboard')
    
    # Obtener empresa si se especifica
    company = None
    if company_id:
        company = get_object_or_404(Company, id=company_id)
        base_queryset = Employee.objects.filter(company=company)
    else:
        if request.user.is_superuser or request.user.role == CustomUser.Role.SUPERUSER:
            base_queryset = Employee.objects.all()
        else:
            companies = request.user.managed_companies.all()
            base_queryset = Employee.objects.filter(company__in=companies)
    
    # Búsqueda
    # Validar y sanitizar búsqueda
    search_query = sanitize_string(request.GET.get('search', ''), max_length=100)
    if search_query:
        # Escapar caracteres especiales para LIKE
        search_pattern = sanitize_sql_like_pattern(search_query)
        base_queryset = base_queryset.filter(
            Q(first_name__icontains=search_pattern) |
            Q(last_name__icontains=search_pattern) |
            Q(identification__icontains=search_pattern) |
            Q(user__email__icontains=search_pattern) |
            Q(position__icontains=search_pattern) |
            Q(area__icontains=search_pattern)
        )
    
    # Filtros - Validar valores permitidos
    status_filter = request.GET.get('status', '')
    if status_filter not in ['', 'active', 'inactive']:
        status_filter = ''
    if status_filter == 'active':
        base_queryset = base_queryset.filter(user__is_active=True)
    elif status_filter == 'inactive':
        base_queryset = base_queryset.filter(user__is_active=False)
    
    # Validar company_filter de forma segura
    company_filter = request.GET.get('company_filter', '')
    if company_filter and not company_id:
        try:
            company_filter_id = validate_integer(company_filter, min_value=1)
            base_queryset = base_queryset.filter(company_id=company_filter_id)
        except (ValueError, TypeError, ValidationError):
            # Si no es un número válido, ignorar el filtro
            company_filter = ''
    
    # Sanitizar filtro de área
    area_filter = sanitize_string(request.GET.get('area', ''), max_length=100)
    if area_filter:
        area_pattern = sanitize_sql_like_pattern(area_filter)
        base_queryset = base_queryset.filter(area__icontains=area_pattern)
    
    # Ordenamiento - Validar campos permitidos
    order_by = request.GET.get('order_by', 'last_name')
    order_direction = request.GET.get('order_direction', 'asc')
    
    valid_order_fields = ['last_name', 'first_name', 'identification', 'position', 'area', 'created_at']
    if order_by not in valid_order_fields:
        order_by = 'last_name'
    
    if order_direction not in ['asc', 'desc']:
        order_direction = 'asc'
    
    if order_direction == 'desc':
        order_by = f'-{order_by}'
    
    employees = base_queryset.select_related('user', 'company').order_by(order_by)
    
    # Paginación - Validar valores permitidos
    try:
        per_page = validate_integer(request.GET.get('per_page', 20), min_value=1, max_value=200)
    except (ValueError, TypeError, ValidationError):
        per_page = 20
    
    valid_per_page = [10, 20, 30, 50, 100, 200]
    if per_page not in valid_per_page:
        per_page = 20
    
    paginator = Paginator(employees, per_page)
    page_number = request.GET.get('page', 1)
    page_obj = paginator.get_page(page_number)
    
    # Obtener empresas para filtro (si no hay company_id)
    companies_for_filter = []
    if not company_id:
        if request.user.is_superuser or request.user.role == CustomUser.Role.SUPERUSER:
            companies_for_filter = Company.objects.all().order_by('name')
        else:
            companies_for_filter = request.user.managed_companies.all().order_by('name')
    
    # Obtener áreas únicas para filtro
    areas = Employee.objects.values_list('area', flat=True).distinct().order_by('area')
    
    context = {
        'employees': page_obj,
        'company': company,
        'search_query': search_query,
        'status_filter': status_filter,
        'company_filter': company_filter,
        'area_filter': area_filter,
        'order_by': order_by.lstrip('-'),
        'order_direction': order_direction,
        'per_page': per_page,
        'companies_for_filter': companies_for_filter,
        'areas': areas,
        'total_count': paginator.count,
    }
    
    return render(request, 'admin/employee_list.html', context)


@login_required
def evaluation_results(request, company_id=None):
    """Resultados de evaluaciones (agregados y anónimos) con recomendaciones empresariales"""
    if not user_is_admin(request.user):
        messages.error(request, 'No tienes permisos para acceder a esta página')
        return redirect('employee_dashboard')
    
    # Importar servicio de recomendaciones empresariales
    from apps.application.services.employer_recommendations_service import EmployerRecommendationsService
    employer_recs_service = EmployerRecommendationsService()
    
    # Filtrar por empresa si se especifica
    company = None
    if company_id:
        company = get_object_or_404(Company, id=company_id)
        evaluations = Evaluation.objects.filter(
            employee__company=company,
            status=Evaluation.Status.COMPLETED
        )
    else:
        if request.user.is_superuser or request.user.role == CustomUser.Role.SUPERUSER:
            evaluations = Evaluation.objects.filter(status=Evaluation.Status.COMPLETED)
            # No asignar company para superusuarios (vista global)
        else:
            # Para administradores de empresa, obtener sus empresas
            companies = request.user.managed_companies.all()
            # Si solo gestiona una empresa, asignarla automáticamente
            if companies.count() == 1:
                company = companies.first()
            evaluations = Evaluation.objects.filter(
                employee__company__in=companies,
                status=Evaluation.Status.COMPLETED
            )
    
    # Obtener año del filtro
    # Validar año de forma segura
    try:
        year_str = request.GET.get('year', str(datetime.now().year))
        year = validate_year(year_str)
        year_str = str(year)
    except (ValueError, TypeError, ValidationError):
        year = datetime.now().year
        year_str = str(year)
    try:
        year = int(year_str)
    except (ValueError, TypeError):
        year = datetime.now().year
    evaluations = evaluations.filter(year=year)
    
    # Calcular estadísticas agregadas por dimensión
    risk_results = RiskResult.objects.filter(evaluation__in=evaluations).select_related('dimension')
    
    # Obtener todas las dimensiones ordenadas
    from apps.infrastructure.models import Dimension
    all_dimensions = Dimension.objects.all().order_by('order')
    
    # Contar total de evaluaciones completadas
    total_evaluations_count = evaluations.count()
    
    # Agrupar por dimensión y nivel de riesgo
    from django.db.models import Count
    dimension_stats = {}
    
    # Inicializar todas las dimensiones
    for dim in all_dimensions:
        dimension_stats[dim.name] = {
            'dimension': dim,
            'BAJO': 0,
            'MEDIO': 0,
            'ALTO': 0,
            'total': 0,
            'porcentaje_bajo': 0.0,
            'porcentaje_medio': 0.0,
            'porcentaje_alto': 0.0,
        }
    
    # Contar resultados por dimensión y nivel de riesgo
    for result in risk_results:
        dim_name = result.dimension.name
        if dim_name in dimension_stats:
            dimension_stats[dim_name][result.risk_level] += 1
            dimension_stats[dim_name]['total'] += 1
    
    # Calcular porcentajes para cada dimensión
    # El porcentaje se calcula sobre el total de evaluaciones completadas
    # Cada evaluación tiene un resultado por dimensión, por lo que el total de evaluaciones
    # es el denominador correcto para calcular el porcentaje de personas en cada nivel de riesgo
    for dim_name, stats in dimension_stats.items():
        if total_evaluations_count > 0:
            stats['porcentaje_bajo'] = round((stats['BAJO'] / total_evaluations_count) * 100, 2)
            stats['porcentaje_medio'] = round((stats['MEDIO'] / total_evaluations_count) * 100, 2)
            stats['porcentaje_alto'] = round((stats['ALTO'] / total_evaluations_count) * 100, 2)
    
    # Crear lista ordenada de dimensiones con estadísticas para el template
    dimension_stats_list = [
        {
            'name': dim.name,
            'order': dim.order,
            'stats': dimension_stats[dim.name]
        }
        for dim in all_dimensions
        if dim.name in dimension_stats
    ]
    
    # ==================== ANÁLISIS DEMOGRÁFICO ====================
    # Función auxiliar para calcular estadísticas por campo demográfico
    def calculate_demographic_stats(field_name, display_method=None):
        """
        Calcula estadísticas de riesgo agrupadas por un campo demográfico
        
        Args:
            field_name: Nombre del campo en el modelo Evaluation
            display_method: Método para obtener el display label (e.g., 'get_gender_display')
        
        Returns:
            Lista de diccionarios con estadísticas por cada valor del campo
        """
        from django.db.models import Count, Q
        
        # Obtener valores únicos del campo
        field_values = evaluations.values_list(field_name, flat=True).distinct().order_by(field_name)
        
        stats = []
        for value in field_values:
            if not value:  # Saltar valores vacíos
                continue
            
            # Filtrar evaluaciones por este valor
            filtered_evals = evaluations.filter(**{field_name: value})
            count = filtered_evals.count()
            
            # Obtener label legible
            if display_method:
                # Obtener una evaluación de ejemplo para usar get_FOO_display()
                sample_eval = filtered_evals.first()
                if sample_eval:
                    label = getattr(sample_eval, display_method)()
                else:
                    label = value
            else:
                label = value
            
            # Calcular riesgo por dimensión para este grupo
            group_risk_results = RiskResult.objects.filter(
                evaluation__in=filtered_evals
            ).select_related('dimension')
            
            # Inicializar contadores por dimensión
            dimension_risks = {}
            for dim in all_dimensions:
                dimension_risks[dim.name] = {
                    'BAJO': 0,
                    'MEDIO': 0,
                    'ALTO': 0,
                    'porcentaje_bajo': 0.0,
                    'porcentaje_medio': 0.0,
                    'porcentaje_alto': 0.0,
                }
            
            # Contar resultados
            for result in group_risk_results:
                dim_name = result.dimension.name
                if dim_name in dimension_risks:
                    dimension_risks[dim_name][result.risk_level] += 1
            
            # Calcular porcentajes
            if count > 0:
                for dim_name, risks in dimension_risks.items():
                    risks['porcentaje_bajo'] = round((risks['BAJO'] / count) * 100, 2)
                    risks['porcentaje_medio'] = round((risks['MEDIO'] / count) * 100, 2)
                    risks['porcentaje_alto'] = round((risks['ALTO'] / count) * 100, 2)
            
            # Generar recomendaciones para dimensiones con riesgo medio y alto
            group_recommendations = {}
            for dim in all_dimensions:
                dim_name = dim.name
                if dim_name in dimension_risks:
                    risks = dimension_risks[dim_name]
                    # Identificar el nivel de riesgo predominante
                    if risks['porcentaje_alto'] > risks['porcentaje_medio'] and risks['porcentaje_alto'] > risks['porcentaje_bajo']:
                        # Riesgo alto predominante
                        if risks['ALTO'] > 0:  # Solo si hay personas en alto riesgo
                            group_recommendations[dim_name] = {
                                'level': 'ALTO',
                                'recommendations': employer_recs_service.get_recommendations_by_dimension(dim_name, 'ALTO'),
                                'count': risks['ALTO'],
                                'percentage': risks['porcentaje_alto']
                            }
                    elif risks['porcentaje_medio'] >= risks['porcentaje_bajo'] and risks['porcentaje_medio'] >= risks['porcentaje_alto']:
                        # Riesgo medio predominante
                        if risks['MEDIO'] > 0:  # Solo si hay personas en medio riesgo
                            group_recommendations[dim_name] = {
                                'level': 'MEDIO',
                                'recommendations': employer_recs_service.get_recommendations_by_dimension(dim_name, 'MEDIO'),
                                'count': risks['MEDIO'],
                                'percentage': risks['porcentaje_medio']
                            }
            
            stats.append({
                'value': value,
                'label': label,
                'count': count,
                'dimension_risks': dimension_risks,
                'recommendations': group_recommendations
            })
        
        return stats
    
    # ==================== ANÁLISIS DE CUMPLIMIENTO ====================
    # Obtener todos los empleados programados (con año de la evaluación)
    from apps.infrastructure.models import Employee
    
    if company_id:
        programmed_employees = Employee.objects.filter(company=company)
    else:
        if request.user.is_superuser or request.user.role == CustomUser.Role.SUPERUSER:
            programmed_employees = Employee.objects.all()
        else:
            companies = request.user.managed_companies.all()
            programmed_employees = Employee.objects.filter(company__in=companies)
    
    total_programmed = programmed_employees.count()
    total_completed = total_evaluations_count
    # Asegurar que el porcentaje sea un número válido entre 0 y 100
    if total_programmed > 0:
        completion_percentage = min(100.0, max(0.0, round((total_completed / total_programmed * 100), 2)))
    else:
        completion_percentage = 0.0
    # Convertir a float para asegurar que se renderice correctamente
    completion_percentage = float(completion_percentage)
    # Asegurar que no exceda 100 ni sea menor a 0
    if completion_percentage > 100:
        completion_percentage = 100.0
    if completion_percentage < 0:
        completion_percentage = 0.0
    total_pending = total_programmed - total_completed
    
    # Función para calcular cumplimiento por categoría demográfica
    def calculate_completion_stats(employee_field, evaluation_field=None, display_method=None):
        """
        Calcula estadísticas de cumplimiento por categoría demográfica.
        Compara empleados programados vs evaluaciones completadas.
        
        Args:
            employee_field: nombre del campo en el modelo Employee
            evaluation_field: nombre del campo en el modelo Evaluation (si es diferente)
            display_method: método para obtener el display label
        """
        # Si no se especifica evaluation_field, usar el mismo nombre
        if evaluation_field is None:
            evaluation_field = employee_field
        
        stats = []
        
        # Obtener valores únicos del campo en empleados programados
        field_values = programmed_employees.values_list(employee_field, flat=True).distinct().order_by(employee_field)
        
        for value in field_values:
            if not value:  # Saltar valores vacíos
                continue
            
            # Contar empleados programados en esta categoría
            programmed_in_category = programmed_employees.filter(**{employee_field: value}).count()
            
            # Contar evaluaciones completadas en esta categoría
            # Usar el campo correspondiente del modelo Evaluation
            completed_in_category = evaluations.filter(**{evaluation_field: value}).count()
            
            # Calcular cumplimiento
            completion_rate = round((completed_in_category / programmed_in_category * 100), 2) if programmed_in_category > 0 else 0
            pending_in_category = programmed_in_category - completed_in_category
            
            # Obtener label legible
            if display_method:
                # Intentar obtener de Evaluation primero (tiene los choices definidos)
                sample_eval = evaluations.filter(**{evaluation_field: value}).first()
                if sample_eval and hasattr(sample_eval, display_method):
                    label = getattr(sample_eval, display_method)()
                else:
                    # Si no hay evaluación o no tiene el método, usar Employee
                    sample_employee = programmed_employees.filter(**{employee_field: value}).first()
                    if sample_employee and hasattr(sample_employee, display_method):
                        label = getattr(sample_employee, display_method)()
                    else:
                        # Si tampoco tiene el método, usar el valor directo
                        label = value
            else:
                label = value
            
            stats.append({
                'value': value,
                'label': label,
                'programmed': programmed_in_category,
                'completed': completed_in_category,
                'pending': pending_in_category,
                'completion_rate': completion_rate
            })
        
        return stats
    
    # Calcular estadísticas de cumplimiento por cada categoría
    # Nota: age_range y experience_range solo existen en Evaluation, no en Employee
    completion_data = {}
    if total_programmed > 0:
        completion_data = {
            'gender': calculate_completion_stats('gender', None, 'get_gender_display'),
            'work_area': calculate_completion_stats('work_area_erp', 'work_area', 'get_work_area_display'),  # Employee.work_area_erp -> Evaluation.work_area
            'education_level': calculate_completion_stats('education_level', None, 'get_education_level_display'),
            'ethnicity': calculate_completion_stats('ethnicity', None, 'get_ethnicity_display'),
            'province': calculate_completion_stats('province', None, None),
            'city': calculate_completion_stats('city', None, None),
        }
        
        # Ahora age_range y experience_range pueden calcularse desde Employee!
        # Calcular estadísticas usando los métodos age_range y experience_range del modelo Employee
        def calculate_completion_stats_from_computed(field_getter_name, evaluation_field, display_method_name):
            """
            Calcula cumplimiento para campos que son propiedades computadas en Employee
            field_getter_name: nombre de la propiedad en Employee (ej: 'age_range')
            evaluation_field: nombre del campo en Evaluation
            display_method_name: nombre del método para display
            """
            stats_dict = {}
            
            # Obtener todos los valores posibles de empleados programados
            for employee in programmed_employees:
                value = getattr(employee, field_getter_name)
                if not value:
                    continue
                
                if value not in stats_dict:
                    stats_dict[value] = {
                        'value': value,
                        'label': getattr(employee, display_method_name)(),
                        'programmed': 0,
                        'completed': 0
                    }
                stats_dict[value]['programmed'] += 1
            
            # Contar completadas por este campo
            for eval in evaluations:
                value = getattr(eval, evaluation_field)
                if value and value in stats_dict:
                    stats_dict[value]['completed'] += 1
            
            # Calcular pending y completion_rate
            stats = []
            for stat in stats_dict.values():
                stat['pending'] = stat['programmed'] - stat['completed']
                stat['completion_rate'] = round((stat['completed'] / stat['programmed'] * 100), 2) if stat['programmed'] > 0 else 0
                stats.append(stat)
            
            return sorted(stats, key=lambda x: x['value'])
        
        completion_data['age_range'] = calculate_completion_stats_from_computed('age_range', 'age_range', 'get_age_range_display')
        completion_data['experience_range'] = calculate_completion_stats_from_computed('experience_range', 'experience_range', 'get_experience_range_display')
    
    # Calcular estadísticas demográficas
    demographic_data = {}
    
    if total_evaluations_count > 0:
        demographic_data = {
            'gender': calculate_demographic_stats('gender', 'get_gender_display'),
            'age_range': calculate_demographic_stats('age_range', 'get_age_range_display'),
            'work_area': calculate_demographic_stats('work_area', 'get_work_area_display'),
            'education_level': calculate_demographic_stats('education_level', 'get_education_level_display'),
            'experience_range': calculate_demographic_stats('experience_range', 'get_experience_range_display'),
            'ethnicity': calculate_demographic_stats('ethnicity', 'get_ethnicity_display'),
            'province': calculate_demographic_stats('province'),
            'city': calculate_demographic_stats('city'),
        }
    
    # ==================== CÁLCULO DE RIESGO GLOBAL (BASADO EN PUNTAJE TOTAL) ====================
    # Según MDT: Calcular el puntaje total de cada evaluación (suma de 58 preguntas)
    # y clasificar según rangos: Bajo (175-232), Medio (117-174), Alto (58-116)
    
    from apps.application.services.risk_calculator import RiskCalculatorService
    calculator = RiskCalculatorService()
    
    global_risk_evaluations = {'bajo': 0, 'medio': 0, 'alto': 0}
    
    for evaluation in evaluations:
        # Calcular puntaje global de esta evaluación
        global_score = calculator.calculate_global_score(evaluation)
        global_risk_level = calculator.calculate_global_risk_level(global_score)
        
        # Clasificar según nivel
        if global_risk_level == 'BAJO':
            global_risk_evaluations['bajo'] += 1
        elif global_risk_level == 'MEDIO':
            global_risk_evaluations['medio'] += 1
        elif global_risk_level == 'ALTO':
            global_risk_evaluations['alto'] += 1
    
    # Calcular porcentajes globales (basado en evaluaciones, no dimensiones)
    global_risk_percentages = {
        'bajo': round((global_risk_evaluations['bajo'] / total_evaluations_count * 100), 2) if total_evaluations_count > 0 else 0,
        'medio': round((global_risk_evaluations['medio'] / total_evaluations_count * 100), 2) if total_evaluations_count > 0 else 0,
        'alto': round((global_risk_evaluations['alto'] / total_evaluations_count * 100), 2) if total_evaluations_count > 0 else 0,
    }
    
    global_risk_counts = global_risk_evaluations  # Alias para compatibilidad
    
    # ==================== ANÁLISIS POR DIMENSIONES (COMPLEMENTARIO) ====================
    # Contar dimensiones predominantes (para análisis adicional)
    dimension_risk_counts = {'bajo': 0, 'medio': 0, 'alto': 0}
    
    if dimension_stats:
        for dim_name, stats in dimension_stats.items():
            bajo_count = stats['BAJO']
            medio_count = stats['MEDIO']
            alto_count = stats['ALTO']
            
            if alto_count > medio_count and alto_count > bajo_count:
                dimension_risk_counts['alto'] += 1
            elif medio_count >= bajo_count and medio_count >= alto_count:
                dimension_risk_counts['medio'] += 1
            else:
                dimension_risk_counts['bajo'] += 1
    
    total_dimensions = len(all_dimensions)
    dimension_risk_percentages = {
        'bajo': round((dimension_risk_counts['bajo'] / total_dimensions * 100), 2) if total_dimensions > 0 else 0,
        'medio': round((dimension_risk_counts['medio'] / total_dimensions * 100), 2) if total_dimensions > 0 else 0,
        'alto': round((dimension_risk_counts['alto'] / total_dimensions * 100), 2) if total_dimensions > 0 else 0,
    }
    
    # Generar recomendaciones empresariales por dimensión y nivel de riesgo
    employer_recommendations = {}
    for dim in all_dimensions:
        employer_recommendations[dim.name] = {
            'BAJO': employer_recs_service.get_recommendations_by_dimension(dim.name, 'BAJO'),
            'MEDIO': employer_recs_service.get_recommendations_by_dimension(dim.name, 'MEDIO'),
            'ALTO': employer_recs_service.get_recommendations_by_dimension(dim.name, 'ALTO'),
        }
    
    # Determinar nivel de riesgo general predominante
    if global_risk_counts['alto'] > global_risk_counts['medio'] and global_risk_counts['alto'] > global_risk_counts['bajo']:
        overall_risk = 'ALTO'
    elif global_risk_counts['medio'] >= global_risk_counts['bajo'] and global_risk_counts['medio'] >= global_risk_counts['alto']:
        overall_risk = 'MEDIO'
    else:
        overall_risk = 'BAJO'
    
    # Obtener recomendaciones generales
    general_recommendations = employer_recs_service.get_general_recommendations_by_level(overall_risk)
    
    context = {
        'company': company,
        'year': year,
        'total_evaluations': total_evaluations_count,
        'dimension_stats': dimension_stats,
        'dimension_stats_list': dimension_stats_list,
        'demographic_data': demographic_data,
        'all_dimensions': all_dimensions,
        'employer_recommendations': employer_recommendations,  # Recomendaciones empresariales
        'general_recommendations': general_recommendations,  # Recomendaciones generales
        'overall_risk': overall_risk,  # Nivel de riesgo general
        'global_risk_counts': global_risk_counts,  # Conteo de EVALUACIONES por nivel (basado en puntaje total)
        'global_risk_percentages': global_risk_percentages,  # Porcentajes globales (basado en puntaje total)
        'dimension_risk_counts': dimension_risk_counts,  # Conteo de DIMENSIONES por nivel (análisis complementario)
        'dimension_risk_percentages': dimension_risk_percentages,  # Porcentajes de dimensiones
        # Datos de cumplimiento
        'total_programmed': total_programmed,  # Total de empleados programados
        'total_completed': total_completed,  # Total de evaluaciones completadas
        'total_pending': total_pending,  # Total de evaluaciones pendientes
        'completion_percentage': completion_percentage,  # Porcentaje de cumplimiento general
        'completion_data': completion_data,  # Datos de cumplimiento por categoría demográfica
    }
    
    return render(request, 'admin/evaluation_results.html', context)


@login_required
def download_admin_results_pdf(request, company_id=None):
    """Genera y descarga un PDF completo con los resultados agregados de la empresa"""
    if not user_is_admin(request.user):
        messages.error(request, 'No tienes permisos para acceder a esta página')
        return redirect('employee_dashboard')
    
    # Importar reportlab solo cuando se necesite
    try:
        from reportlab.lib.pagesizes import A4, landscape
        from reportlab.lib import colors
        from reportlab.lib.units import inch
        from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer, PageBreak
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.enums import TA_CENTER, TA_JUSTIFY, TA_LEFT
        from io import BytesIO
    except ImportError:
        messages.error(request, 'Error: La librería reportlab no está instalada correctamente.')
        return redirect('evaluation_results')
    
    # Importar servicio de recomendaciones empresariales
    from apps.application.services.employer_recommendations_service import EmployerRecommendationsService
    employer_recs_service = EmployerRecommendationsService()
    
    # Filtrar por empresa si se especifica (reutilizar lógica de evaluation_results)
    company = None
    if company_id:
        company = get_object_or_404(Company, id=company_id)
        evaluations = Evaluation.objects.filter(
            employee__company=company,
            status=Evaluation.Status.COMPLETED
        )
    else:
        if request.user.is_superuser or request.user.role == CustomUser.Role.SUPERUSER:
            evaluations = Evaluation.objects.filter(status=Evaluation.Status.COMPLETED)
        else:
            companies = request.user.managed_companies.all()
            if companies.count() == 1:
                company = companies.first()
            evaluations = Evaluation.objects.filter(
                employee__company__in=companies,
                status=Evaluation.Status.COMPLETED
            )
    
    # Obtener año del filtro
    # Validar año de forma segura
    try:
        year_str = request.GET.get('year', str(datetime.now().year))
        year = validate_year(year_str)
        year_str = str(year)
    except (ValueError, TypeError, ValidationError):
        year = datetime.now().year
        year_str = str(year)
    try:
        year = int(year_str)
    except (ValueError, TypeError):
        year = datetime.now().year
    evaluations = evaluations.filter(year=year)
    
    # Calcular estadísticas (similar a evaluation_results)
    risk_results = RiskResult.objects.filter(evaluation__in=evaluations).select_related('dimension')
    from apps.infrastructure.models import Dimension
    all_dimensions = Dimension.objects.all().order_by('order')
    total_evaluations_count = evaluations.count()
    
    # Calcular dimension_stats
    dimension_stats = {}
    for dim in all_dimensions:
        dimension_stats[dim.name] = {
            'dimension': dim,
            'BAJO': 0,
            'MEDIO': 0,
            'ALTO': 0,
            'total': 0,
            'porcentaje_bajo': 0.0,
            'porcentaje_medio': 0.0,
            'porcentaje_alto': 0.0,
        }
    
    for result in risk_results:
        dim_name = result.dimension.name
        if dim_name in dimension_stats:
            dimension_stats[dim_name][result.risk_level] += 1
            dimension_stats[dim_name]['total'] += 1
    
    for dim_name, stats in dimension_stats.items():
        if total_evaluations_count > 0:
            stats['porcentaje_bajo'] = round((stats['BAJO'] / total_evaluations_count) * 100, 2)
            stats['porcentaje_medio'] = round((stats['MEDIO'] / total_evaluations_count) * 100, 2)
            stats['porcentaje_alto'] = round((stats['ALTO'] / total_evaluations_count) * 100, 2)
    
    # ==================== CÁLCULO DE RIESGO GLOBAL (BASADO EN PUNTAJE TOTAL) ====================
    from apps.application.services.risk_calculator import RiskCalculatorService
    calculator = RiskCalculatorService()
    
    global_risk_evaluations = {'bajo': 0, 'medio': 0, 'alto': 0}
    
    for evaluation in evaluations:
        global_score = calculator.calculate_global_score(evaluation)
        global_risk_level = calculator.calculate_global_risk_level(global_score)
        
        if global_risk_level == 'BAJO':
            global_risk_evaluations['bajo'] += 1
        elif global_risk_level == 'MEDIO':
            global_risk_evaluations['medio'] += 1
        elif global_risk_level == 'ALTO':
            global_risk_evaluations['alto'] += 1
    
    global_risk_percentages = {
        'bajo': round((global_risk_evaluations['bajo'] / total_evaluations_count * 100), 1) if total_evaluations_count > 0 else 0,
        'medio': round((global_risk_evaluations['medio'] / total_evaluations_count * 100), 1) if total_evaluations_count > 0 else 0,
        'alto': round((global_risk_evaluations['alto'] / total_evaluations_count * 100), 1) if total_evaluations_count > 0 else 0,
    }
    
    global_risk_counts = global_risk_evaluations  # Alias
    
    # Determinar nivel de riesgo general predominante (usar la MISMA lógica que evaluation_results)
    if global_risk_counts['alto'] > global_risk_counts['medio'] and global_risk_counts['alto'] > global_risk_counts['bajo']:
        overall_risk = 'ALTO'
    elif global_risk_counts['medio'] >= global_risk_counts['bajo'] and global_risk_counts['medio'] >= global_risk_counts['alto']:
        overall_risk = 'MEDIO'
    else:
        overall_risk = 'BAJO'
    
    # Recomendaciones generales
    general_recommendations = employer_recs_service.get_general_recommendations_by_level(overall_risk)
    
    # ==================== CALCULAR DATOS DE CUMPLIMIENTO ====================
    # Obtener todas las empresas objetivo
    if company:
        target_companies = [company]
    elif request.user.is_superuser or request.user.role == CustomUser.Role.SUPERUSER:
        target_companies = Company.objects.all()
    else:
        target_companies = list(request.user.managed_companies.all())
    
    # Calcular estadísticas de cumplimiento
    total_programmed = Employee.objects.filter(company__in=target_companies).count()
    total_completed = evaluations.count()
    total_pending = total_programmed - total_completed
    overall_compliance_percentage = round((total_completed / total_programmed * 100), 2) if total_programmed > 0 else 0
    
    # Función auxiliar para calcular cumplimiento (SIN DUPLICADOS)
    def calculate_completion_stats_for_pdf(employee_field, evaluation_field=None, display_method=None):
        if evaluation_field is None:
            evaluation_field = employee_field
        
        stats_dict = {}  # Usar diccionario para evitar duplicados
        
        if employee_field in ['age_range', 'experience_range']:
            # Manejo especial para campos calculados
            if employee_field == 'age_range':
                all_ranges = list(set(Employee.AgeRange.values) | set(evaluations.values_list('age_range', flat=True)))
                for ar_value in all_ranges:
                    if not ar_value: continue
                    programmed_count = sum(1 for emp in Employee.objects.filter(company__in=target_companies) if emp.age_range == ar_value)
                    completed_count = evaluations.filter(age_range=ar_value).count()
                    if programmed_count > 0 or completed_count > 0:
                        label = dict(Employee.AgeRange.choices).get(ar_value, ar_value)
                        stats_dict[ar_value] = {
                            'label': label,
                            'programmed': programmed_count,
                            'completed': completed_count,
                            'percentage': round((completed_count / programmed_count * 100), 1) if programmed_count > 0 else 0
                        }
            elif employee_field == 'experience_range':
                all_ranges = list(set(Employee.ExperienceRange.values) | set(evaluations.values_list('experience_range', flat=True)))
                for er_value in all_ranges:
                    if not er_value: continue
                    programmed_count = sum(1 for emp in Employee.objects.filter(company__in=target_companies) if emp.experience_range == er_value)
                    completed_count = evaluations.filter(experience_range=er_value).count()
                    if programmed_count > 0 or completed_count > 0:
                        label = dict(Employee.ExperienceRange.choices).get(er_value, er_value)
                        stats_dict[er_value] = {
                            'label': label,
                            'programmed': programmed_count,
                            'completed': completed_count,
                            'percentage': round((completed_count / programmed_count * 100), 1) if programmed_count > 0 else 0
                        }
        else:
            # Campos normales - USAR .distinct() correctamente
            employee_values = Employee.objects.filter(company__in=target_companies).values_list(employee_field, flat=True).distinct()
            unique_values = set(employee_values)  # Eliminar duplicados definitivamente
            
            for value in unique_values:
                if not value: continue
                
                programmed_count = Employee.objects.filter(company__in=target_companies, **{employee_field: value}).count()
                completed_count = evaluations.filter(**{evaluation_field: value}).count()
                
                label = value
                if display_method:
                    sample_eval = evaluations.filter(**{evaluation_field: value}).first()
                    if sample_eval and hasattr(sample_eval, display_method):
                        label = getattr(sample_eval, display_method)()
                    else:
                        sample_emp = Employee.objects.filter(company__in=target_companies, **{employee_field: value}).first()
                        if sample_emp and hasattr(sample_emp, display_method):
                            label = getattr(sample_emp, display_method)()
                
                if programmed_count > 0 or completed_count > 0:
                    stats_dict[value] = {
                        'label': label,
                        'programmed': programmed_count,
                        'completed': completed_count,
                        'percentage': round((completed_count / programmed_count * 100), 1) if programmed_count > 0 else 0
                    }
        
        # Convertir diccionario a lista (sin duplicados)
        return list(stats_dict.values())
    
    # Calcular cumplimiento por categoría
    compliance_by_gender = calculate_completion_stats_for_pdf('gender', None, 'get_gender_display')
    compliance_by_area = calculate_completion_stats_for_pdf('work_area_erp', 'work_area', 'get_work_area_display')
    compliance_by_education = calculate_completion_stats_for_pdf('education_level', None, 'get_education_level_display')
    
    # ==================== CALCULAR DATOS DEMOGRÁFICOS ====================
    def calculate_demographic_stats_for_pdf(field_name, display_method=None):
        """Calcula estadísticas demográficas agrupadas"""
        from django.db.models import Count, Q
        
        groups = evaluations.values(field_name).annotate(count=Count('id'))
        result = []
        
        for group in groups:
            value = group[field_name]
            if not value:
                continue
            
            count = group['count']
            group_evals = evaluations.filter(**{field_name: value})
            
            # Calcular dimension_risks para este grupo
            dimension_risks = {}
            for dim in all_dimensions:
                group_results = RiskResult.objects.filter(
                    evaluation__in=group_evals,
                    dimension=dim
                )
                bajo = group_results.filter(risk_level='BAJO').count()
                medio = group_results.filter(risk_level='MEDIO').count()
                alto = group_results.filter(risk_level='ALTO').count()
                total = bajo + medio + alto
                
                dimension_risks[dim.name] = {
                    'bajo': bajo,
                    'medio': medio,
                    'alto': alto,
                    'porcentaje_bajo': round((bajo / total * 100), 1) if total > 0 else 0,
                    'porcentaje_medio': round((medio / total * 100), 1) if total > 0 else 0,
                    'porcentaje_alto': round((alto / total * 100), 1) if total > 0 else 0,
                }
            
            # Obtener label
            label = value
            if display_method:
                sample = group_evals.first()
                if sample and hasattr(sample, display_method):
                    label = getattr(sample, display_method)()
            
            result.append({
                'value': value,
                'label': label,
                'count': count,
                'dimension_risks': dimension_risks
            })
        
        return result
    
    demographic_data_pdf = {
        'gender': calculate_demographic_stats_for_pdf('gender', 'get_gender_display'),
        'age_range': calculate_demographic_stats_for_pdf('age_range', 'get_age_range_display'),
        'work_area': calculate_demographic_stats_for_pdf('work_area', 'get_work_area_display'),
        'education_level': calculate_demographic_stats_for_pdf('education_level', 'get_education_level_display'),
    }
    
    # Crear el PDF con márgenes reducidos
    buffer = BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=A4, rightMargin=30, leftMargin=30, topMargin=30, bottomMargin=30)
    elements = []
    
    # Estilos
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(
        'CustomTitle',
        parent=styles['Heading1'],
        fontSize=16,
        textColor=colors.HexColor('#1e40af'),
        spaceAfter=12,
        alignment=TA_CENTER,
        fontName='Helvetica-Bold'
    )
    
    heading_style = ParagraphStyle(
        'CustomHeading',
        parent=styles['Heading2'],
        fontSize=12,
        textColor=colors.HexColor('#1e3a8a'),
        spaceAfter=6,
        spaceBefore=8,
        fontName='Helvetica-Bold'
    )
    
    subheading_style = ParagraphStyle(
        'CustomSubHeading',
        parent=styles['Heading3'],
        fontSize=10,
        textColor=colors.HexColor('#334155'),
        spaceAfter=4,
        spaceBefore=6,
        fontName='Helvetica-Bold'
    )
    
    normal_style = ParagraphStyle(
        'CustomNormal',
        parent=styles['Normal'],
        fontSize=8,
        textColor=colors.HexColor('#1f2937'),
        spaceAfter=3,
        alignment=TA_JUSTIFY,
        leading=10
    )
    
    # ==================== PORTADA ====================
    company_name = company.name if company else "Resultados Globales"
    title = Paragraph(f"Informe de Evaluación de Riesgo Psicosocial<br/>{company_name}<br/>Año {year}", title_style)
    elements.append(title)
    elements.append(Spacer(1, 0.15*inch))
    
    # Información general
    info_data = [
        ['Empresa:', company_name],
        ['Año de Evaluación:', str(year)],
        ['Total de Evaluaciones:', str(total_evaluations_count)],
        ['Fecha del Informe:', datetime.now().strftime('%d/%m/%Y')],
    ]
    
    info_table = Table(info_data, colWidths=[2*inch, 3.5*inch])
    info_table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (0, -1), colors.HexColor('#e0e7ff')),
        ('TEXTCOLOR', (0, 0), (-1, -1), colors.HexColor('#1e293b')),
        ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
        ('FONTNAME', (0, 0), (0, -1), 'Helvetica-Bold'),
        ('FONTNAME', (1, 0), (1, -1), 'Helvetica'),
        ('FONTSIZE', (0, 0), (-1, -1), 9),
        ('BOTTOMPADDING', (0, 0), (-1, -1), 4),
        ('TOPPADDING', (0, 0), (-1, -1), 4),
        ('GRID', (0, 0), (-1, -1), 1, colors.HexColor('#cbd5e1')),
    ]))
    elements.append(info_table)
    elements.append(Spacer(1, 0.2*inch))
    
    # ==================== SECCIÓN 1: CUMPLIMIENTO ====================
    compliance_title = Paragraph("1. ANÁLISIS DE CUMPLIMIENTO", heading_style)
    elements.append(compliance_title)
    elements.append(Spacer(1, 0.05*inch))
    
    # Resumen general de cumplimiento
    compliance_summary_data = [
        ['Concepto', 'Cantidad', 'Porcentaje'],
        ['Empleados Programados', str(total_programmed), '100%'],
        ['Evaluaciones Completadas', str(total_completed), f'{overall_compliance_percentage:.1f}%'],
        ['Evaluaciones Pendientes', str(total_pending), f'{100 - overall_compliance_percentage:.1f}%'],
    ]
    
    compliance_summary_table = Table(compliance_summary_data, colWidths=[2.5*inch, 1.5*inch, 1.5*inch])
    compliance_summary_table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#3b82f6')),
        ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
        ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
        ('FONTSIZE', (0, 0), (-1, 0), 9),
        ('BOTTOMPADDING', (0, 0), (-1, 0), 5),
        ('TOPPADDING', (0, 0), (-1, 0), 5),
        ('FONTNAME', (0, 1), (-1, -1), 'Helvetica'),
        ('FONTSIZE', (0, 1), (-1, -1), 8),
        ('BOTTOMPADDING', (0, 1), (-1, -1), 3),
        ('TOPPADDING', (0, 1), (-1, -1), 3),
        ('GRID', (0, 0), (-1, -1), 1, colors.HexColor('#cbd5e1')),
        ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.HexColor('#f8fafc')]),
    ]))
    elements.append(compliance_summary_table)
    elements.append(Spacer(1, 0.1*inch))
    
    # Cumplimiento por Género
    if compliance_by_gender:
        gender_heading = Paragraph("Cumplimiento por Género", subheading_style)
        elements.append(gender_heading)
        
        gender_data = [['Género', 'Programados', 'Completados', 'Cumplimiento']]
        for stat in compliance_by_gender:
            gender_data.append([
                stat['label'],
                str(stat['programmed']),
                str(stat['completed']),
                f"{stat['percentage']:.1f}%"
            ])
        
        gender_table = Table(gender_data, colWidths=[1.8*inch, 1.3*inch, 1.3*inch, 1.3*inch])
        gender_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#8b5cf6')),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 8),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 4),
            ('TOPPADDING', (0, 0), (-1, 0), 4),
            ('FONTNAME', (0, 1), (-1, -1), 'Helvetica'),
            ('FONTSIZE', (0, 1), (-1, -1), 7),
            ('BOTTOMPADDING', (0, 1), (-1, -1), 3),
            ('TOPPADDING', (0, 1), (-1, -1), 3),
            ('GRID', (0, 0), (-1, -1), 1, colors.HexColor('#cbd5e1')),
        ]))
        elements.append(gender_table)
        elements.append(Spacer(1, 0.08*inch))
    
    # Cumplimiento por Área
    if compliance_by_area:
        area_heading = Paragraph("Cumplimiento por Área de Trabajo", subheading_style)
        elements.append(area_heading)
        
        area_data = [['Área', 'Programados', 'Completados', 'Cumplimiento']]
        for stat in compliance_by_area:
            area_data.append([
                Paragraph(stat['label'], normal_style),
                str(stat['programmed']),
                str(stat['completed']),
                f"{stat['percentage']:.1f}%"
            ])
        
        area_table = Table(area_data, colWidths=[2.2*inch, 1.2*inch, 1.2*inch, 1.1*inch])
        area_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#14b8a6')),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
            ('ALIGN', (0, 0), (0, -1), 'LEFT'),
            ('ALIGN', (1, 0), (-1, -1), 'CENTER'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 8),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 4),
            ('TOPPADDING', (0, 0), (-1, 0), 4),
            ('FONTNAME', (0, 1), (-1, -1), 'Helvetica'),
            ('FONTSIZE', (0, 1), (-1, -1), 7),
            ('BOTTOMPADDING', (0, 1), (-1, -1), 3),
            ('TOPPADDING', (0, 1), (-1, -1), 3),
            ('GRID', (0, 0), (-1, -1), 1, colors.HexColor('#cbd5e1')),
        ]))
        elements.append(area_table)
        elements.append(Spacer(1, 0.08*inch))
    
    # Cumplimiento por Educación
    if compliance_by_education:
        edu_heading = Paragraph("Cumplimiento por Nivel de Educación", subheading_style)
        elements.append(edu_heading)
        
        edu_data = [['Nivel Educativo', 'Programados', 'Completados', 'Cumplimiento']]
        for stat in compliance_by_education:
            edu_data.append([
                stat['label'],
                str(stat['programmed']),
                str(stat['completed']),
                f"{stat['percentage']:.1f}%"
            ])
        
        edu_table = Table(edu_data, colWidths=[2*inch, 1.3*inch, 1.3*inch, 1.1*inch])
        edu_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#f59e0b')),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 8),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 4),
            ('TOPPADDING', (0, 0), (-1, 0), 4),
            ('FONTNAME', (0, 1), (-1, -1), 'Helvetica'),
            ('FONTSIZE', (0, 1), (-1, -1), 7),
            ('BOTTOMPADDING', (0, 1), (-1, -1), 3),
            ('TOPPADDING', (0, 1), (-1, -1), 3),
            ('GRID', (0, 0), (-1, -1), 1, colors.HexColor('#cbd5e1')),
        ]))
        elements.append(edu_table)
    
    elements.append(PageBreak())
    
    # ==================== SECCIÓN 2: RESULTADO GLOBAL (VISTA GENERAL) ====================
    vista_general_title = Paragraph("2. VISTA GENERAL - RESULTADOS DE EVALUACIÓN", heading_style)
    elements.append(vista_general_title)
    elements.append(Spacer(1, 0.05*inch))
    global_heading = Paragraph("2.1 Resultado Global de la Evaluación", subheading_style)
    elements.append(global_heading)
    
    risk_color = colors.HexColor('#10b981') if overall_risk == 'BAJO' else (colors.HexColor('#eab308') if overall_risk == 'MEDIO' else colors.HexColor('#ef4444'))
    
    global_data = [
        ['Nivel de Riesgo General (MDT):', overall_risk],
        ['Personas en Riesgo Bajo (175-232 pts):', f"{global_risk_counts['bajo']} persona{'s' if global_risk_counts['bajo'] != 1 else ''} ({global_risk_percentages['bajo']:.1f}%)"],
        ['Personas en Riesgo Medio (117-174 pts):', f"{global_risk_counts['medio']} persona{'s' if global_risk_counts['medio'] != 1 else ''} ({global_risk_percentages['medio']:.1f}%)"],
        ['Personas en Riesgo Alto (58-116 pts):', f"{global_risk_counts['alto']} persona{'s' if global_risk_counts['alto'] != 1 else ''} ({global_risk_percentages['alto']:.1f}%)"],
    ]
    
    global_table = Table(global_data, colWidths=[2.5*inch, 3*inch])
    global_table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (0, -1), colors.HexColor('#f1f5f9')),
        ('TEXTCOLOR', (0, 0), (-1, -1), colors.HexColor('#1e293b')),
        ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
        ('FONTNAME', (0, 0), (0, -1), 'Helvetica-Bold'),
        ('FONTNAME', (1, 0), (1, -1), 'Helvetica'),
        ('FONTSIZE', (0, 0), (-1, -1), 8),
        ('BOTTOMPADDING', (0, 0), (-1, -1), 4),
        ('TOPPADDING', (0, 0), (-1, -1), 4),
        ('BACKGROUND', (1, 0), (1, 0), risk_color),
        ('TEXTCOLOR', (1, 0), (1, 0), colors.white),
        ('GRID', (0, 0), (-1, -1), 1, colors.HexColor('#cbd5e1')),
    ]))
    elements.append(global_table)
    elements.append(Spacer(1, 0.1*inch))
    
    # ==================== RESULTADOS POR DIMENSIÓN ====================
    dimensions_heading = Paragraph("2.2 Resultados por Dimensión", subheading_style)
    elements.append(dimensions_heading)
    
    # Encabezado de la tabla
    dimensions_header = [['Dimensión', 'Bajo %', 'Medio %', 'Alto %', 'Total']]
    
    # Datos de dimensiones
    dimensions_data = []
    for dim in all_dimensions:
        if dim.name in dimension_stats:
            stats = dimension_stats[dim.name]
            dimensions_data.append([
                Paragraph(dim.name, normal_style),
                f"{stats['porcentaje_bajo']:.1f}%\n({stats['BAJO']})",
                f"{stats['porcentaje_medio']:.1f}%\n({stats['MEDIO']})",
                f"{stats['porcentaje_alto']:.1f}%\n({stats['ALTO']})",
                str(stats['total'])
            ])
    
    # Crear tabla completa
    dimensions_table_data = dimensions_header + dimensions_data
    dimensions_table = Table(dimensions_table_data, colWidths=[2.2*inch, 0.9*inch, 0.9*inch, 0.9*inch, 0.6*inch])
    
    # Estilos de la tabla
    table_style = [
        ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#1e40af')),
        ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
        ('ALIGN', (0, 0), (0, -1), 'LEFT'),
        ('ALIGN', (1, 0), (-1, -1), 'CENTER'),
        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
        ('FONTSIZE', (0, 0), (-1, 0), 8),
        ('BOTTOMPADDING', (0, 0), (-1, 0), 4),
        ('TOPPADDING', (0, 0), (-1, 0), 4),
        ('FONTNAME', (0, 1), (-1, -1), 'Helvetica'),
        ('FONTSIZE', (0, 1), (-1, -1), 7),
        ('BOTTOMPADDING', (0, 1), (-1, -1), 3),
        ('TOPPADDING', (0, 1), (-1, -1), 3),
        ('GRID', (0, 0), (-1, -1), 1, colors.HexColor('#cbd5e1')),
        ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.HexColor('#f8fafc')]),
    ]
    dimensions_table.setStyle(TableStyle(table_style))
    elements.append(dimensions_table)
    elements.append(PageBreak())
    
    # ==================== RECOMENDACIONES ESTRATÉGICAS GENERALES ====================
    recs_heading = Paragraph("2.3 Recomendaciones Estratégicas Generales", subheading_style)
    elements.append(recs_heading)
    
    for i, rec in enumerate(general_recommendations[:8], 1):  # Limitar a 8 para no saturar
        rec_para = Paragraph(f"{i}. {rec}", normal_style)
        elements.append(rec_para)
    
    elements.append(Spacer(1, 0.1*inch))
    
    # ==================== RECOMENDACIONES POR DIMENSIÓN (RIESGO MEDIO Y ALTO) ====================
    dim_recs_heading = Paragraph("2.4 Plan de Acción por Dimensión (Riesgo Medio y Alto)", subheading_style)
    elements.append(dim_recs_heading)
    
    for dim in all_dimensions:
        if dim.name in dimension_stats:
            stats = dimension_stats[dim.name]
            # Determinar si mostrar recomendaciones (si hay medio o alto)
            if stats['MEDIO'] > 0 or stats['ALTO'] > 0:
                dim_subheading = Paragraph(f"• {dim.name}", subheading_style)
                elements.append(dim_subheading)
                
                # Determinar nivel predominante
                if stats['porcentaje_alto'] > stats['porcentaje_medio']:
                    level = 'ALTO'
                    count = stats['ALTO']
                    percentage = stats['porcentaje_alto']
                else:
                    level = 'MEDIO'
                    count = stats['MEDIO']
                    percentage = stats['porcentaje_medio']
                
                level_info = Paragraph(f"Riesgo {level}: {percentage:.1f}% ({count} persona{'s' if count != 1 else ''})", normal_style)
                elements.append(level_info)
                
                recs = employer_recs_service.get_recommendations_by_dimension(dim.name, level)
                for rec in recs[:3]:  # Limitar a 3 recomendaciones por dimensión
                    rec_bullet = Paragraph(f"  - {rec}", normal_style)
                    elements.append(rec_bullet)
                
                elements.append(Spacer(1, 0.05*inch))
    
    elements.append(PageBreak())
    
    # ==================== SECCIÓN 3: ANÁLISIS DEMOGRÁFICO ====================
    demographic_title = Paragraph("3. ANÁLISIS DEMOGRÁFICO DETALLADO", heading_style)
    elements.append(demographic_title)
    elements.append(Spacer(1, 0.1*inch))
    
    # 3.1 Análisis por Género
    if demographic_data_pdf.get('gender'):
        gender_demo_heading = Paragraph("3.1 Análisis por Género", subheading_style)
        elements.append(gender_demo_heading)
        
        for group in demographic_data_pdf['gender']:
            group_title = Paragraph(f"<b>{group['label']}</b> ({group['count']} participante{'s' if group['count'] != 1 else ''})", normal_style)
            elements.append(group_title)
            
            # Tabla de dimensiones para este grupo
            group_dim_data = [['Dimensión', 'Bajo %', 'Medio %', 'Alto %']]
            for dim in all_dimensions[:7]:  # Primeras 7 dimensiones
                if dim.name in group['dimension_risks']:
                    risks = group['dimension_risks'][dim.name]
                    group_dim_data.append([
                        Paragraph(dim.name, ParagraphStyle('tiny', parent=normal_style, fontSize=7)),
                        f"{risks['porcentaje_bajo']:.0f}%",
                        f"{risks['porcentaje_medio']:.0f}%",
                        f"{risks['porcentaje_alto']:.0f}%"
                    ])
            
            group_table = Table(group_dim_data, colWidths=[2.5*inch, 0.9*inch, 0.9*inch, 0.9*inch])
            group_table.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#8b5cf6')),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
                ('ALIGN', (0, 0), (0, -1), 'LEFT'),
                ('ALIGN', (1, 0), (-1, -1), 'CENTER'),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('FONTSIZE', (0, 0), (-1, 0), 8),
                ('BOTTOMPADDING', (0, 0), (-1, 0), 6),
                ('TOPPADDING', (0, 0), (-1, 0), 6),
                ('FONTNAME', (0, 1), (-1, -1), 'Helvetica'),
                ('FONTSIZE', (0, 1), (-1, -1), 7),
                ('BOTTOMPADDING', (0, 1), (-1, -1), 4),
                ('TOPPADDING', (0, 1), (-1, -1), 4),
                ('GRID', (0, 0), (-1, -1), 0.5, colors.HexColor('#cbd5e1')),
            ]))
            elements.append(group_table)
            elements.append(Spacer(1, 0.1*inch))
            
            # Recomendaciones para este grupo (solo riesgo medio/alto)
            has_recommendations = False
            for dim in all_dimensions:
                if dim.name in group['dimension_risks']:
                    risks = group['dimension_risks'][dim.name]
                    if risks['medio'] > 0 or risks['alto'] > 0:
                        if not has_recommendations:
                            recs_subheading = Paragraph(f"Recomendaciones específicas para {group['label']}:", 
                                ParagraphStyle('tiny_bold', parent=normal_style, fontSize=8, fontName='Helvetica-Bold'))
                            elements.append(recs_subheading)
                            has_recommendations = True
                        
                        level = 'ALTO' if risks['alto'] > risks['medio'] else 'MEDIO'
                        recs = employer_recs_service.get_recommendations_by_dimension(dim.name, level)
                        if recs:
                            dim_rec_text = Paragraph(f"<b>{dim.name} ({level}):</b> {recs[0][:150]}...", 
                                ParagraphStyle('tiny', parent=normal_style, fontSize=7, leading=9))
                            elements.append(dim_rec_text)
            
            elements.append(Spacer(1, 0.08*inch))
        
        elements.append(PageBreak())
    
    # 3.2 Análisis por Rango de Edad
    if demographic_data_pdf.get('age_range'):
        age_demo_heading = Paragraph("3.2 Análisis por Rango de Edad", subheading_style)
        elements.append(age_demo_heading)
        
        for group in demographic_data_pdf['age_range'][:3]:  # Limitar a 3 grupos
            group_title = Paragraph(f"<b>{group['label']}</b> ({group['count']} participante{'s' if group['count'] != 1 else ''})", normal_style)
            elements.append(group_title)
            
            # Tabla resumida
            summary_data = []
            high_risk_dims = []
            for dim in all_dimensions:
                if dim.name in group['dimension_risks']:
                    risks = group['dimension_risks'][dim.name]
                    if risks['porcentaje_alto'] > 30:  # Solo dimensiones con alto riesgo significativo
                        high_risk_dims.append(f"{dim.name} ({risks['porcentaje_alto']:.0f}% alto)")
            
            if high_risk_dims:
                dims_text = Paragraph(f"<b>Dimensiones de atención:</b> {', '.join(high_risk_dims[:3])}", 
                    ParagraphStyle('tiny', parent=normal_style, fontSize=8))
                elements.append(dims_text)
            
            elements.append(Spacer(1, 0.08*inch))
        
        elements.append(PageBreak())
    
    # 3.3 Análisis por Área de Trabajo
    if demographic_data_pdf.get('work_area'):
        area_demo_heading = Paragraph("3.3 Análisis por Área de Trabajo", subheading_style)
        elements.append(area_demo_heading)
        
        for group in demographic_data_pdf['work_area'][:4]:  # Limitar a 4 áreas
            group_title = Paragraph(f"<b>{group['label']}</b> ({group['count']} participante{'s' if group['count'] != 1 else ''})", normal_style)
            elements.append(group_title)
            
            # Identificar dimensiones críticas
            critical_dims = []
            for dim in all_dimensions:
                if dim.name in group['dimension_risks']:
                    risks = group['dimension_risks'][dim.name]
                    if risks['porcentaje_alto'] > 25 or risks['porcentaje_medio'] > 40:
                        level = 'ALTO' if risks['porcentaje_alto'] > 25 else 'MEDIO'
                        critical_dims.append((dim.name, level, risks))
            
            if critical_dims:
                for dim_name, level, risks in critical_dims[:2]:  # Top 2 críticas
                    recs = employer_recs_service.get_recommendations_by_dimension(dim_name, level)
                    if recs:
                        rec_text = Paragraph(f"• <b>{dim_name}:</b> {recs[0][:120]}...", 
                            ParagraphStyle('tiny', parent=normal_style, fontSize=8, leading=10))
                        elements.append(rec_text)
            
            elements.append(Spacer(1, 0.06*inch))
        
        elements.append(PageBreak())
    
    # 3.4 Análisis por Nivel de Educación
    if demographic_data_pdf.get('education_level'):
        edu_demo_heading = Paragraph("3.4 Análisis por Nivel de Educación", subheading_style)
        elements.append(edu_demo_heading)
        
        for group in demographic_data_pdf['education_level'][:3]:
            group_title = Paragraph(f"<b>{group['label']}</b> ({group['count']} participante{'s' if group['count'] != 1 else ''})", normal_style)
            elements.append(group_title)
            
            # Resumen de riesgos
            total_alto = sum(group['dimension_risks'][d.name]['alto'] for d in all_dimensions if d.name in group['dimension_risks'])
            total_medio = sum(group['dimension_risks'][d.name]['medio'] for d in all_dimensions if d.name in group['dimension_risks'])
            
            summary_para = Paragraph(f"Dimensiones en riesgo alto: {total_alto}, medio: {total_medio}", 
                ParagraphStyle('tiny', parent=normal_style, fontSize=8))
            elements.append(summary_para)
            elements.append(Spacer(1, 0.1*inch))
    
    elements.append(PageBreak())
    
    # ==================== NOTA FINAL ====================
    final_note = Paragraph(
        "<b>Nota Importante:</b> Este informe contiene información agregada y anónima de las evaluaciones "
        "de riesgo psicosocial. Las recomendaciones presentadas deben ser implementadas de manera progresiva "
        "y adaptadas al contexto específico de la organización. Se recomienda establecer un plan de seguimiento "
        "y evaluación periódica de las acciones implementadas.",
        normal_style
    )
    elements.append(final_note)
    
    # Construir el PDF
    doc.build(elements)
    
    # Preparar la respuesta HTTP
    buffer.seek(0)
    response = HttpResponse(buffer, content_type='application/pdf')
    filename = f"Informe_Riesgo_Psicosocial_{company_name.replace(' ', '_')}_{year}.pdf"
    response['Content-Disposition'] = f'attachment; filename="{filename}"'
    
    return response


@login_required
def create_company(request):
    """Crear nueva empresa con usuario admin"""
    if not (request.user.is_superuser or request.user.role == CustomUser.Role.SUPERUSER):
        messages.error(request, 'No tienes permisos para crear empresas')
        return redirect('admin_dashboard')
    
    if request.method == 'POST':
        form = CompanyForm(request.POST)
        if form.is_valid():
            company = form.save()
            messages.success(request, f'Empresa "{company.name}" creada exitosamente')
            messages.info(request, f'Usuario admin creado: {form.cleaned_data["admin_email"]}')
            messages.info(request, f'Contraseña: {form.admin_password}')
            return redirect('company_list')
    else:
        form = CompanyForm()
    
    context = {
        'form': form,
        'title': 'Nueva Empresa',
        'ciudades_por_provincia': json.dumps(CIUDADES_POR_PROVINCIA)
    }
    return render(request, 'admin/company_form.html', context)


@login_required
def edit_company(request, company_id):
    """Editar empresa existente"""
    if not (request.user.is_superuser or request.user.role == CustomUser.Role.SUPERUSER):
        messages.error(request, 'No tienes permisos para editar empresas')
        return redirect('admin_dashboard')
    
    company = get_object_or_404(Company, id=company_id)
    
    if request.method == 'POST':
        form = CompanyForm(request.POST, instance=company)
        if form.is_valid():
            company = form.save()
            messages.success(request, f'Empresa "{company.name}" actualizada exitosamente')
            if form.cleaned_data.get('admin_password'):
                messages.info(request, f'Nueva contraseña: {form.admin_password}')
            return redirect('company_list')
    else:
        form = CompanyForm(instance=company)
    
    context = {
        'form': form,
        'company': company,
        'title': 'Editar Empresa',
        'ciudades_por_provincia': json.dumps(CIUDADES_POR_PROVINCIA)
    }
    return render(request, 'admin/company_form.html', context)


@login_required
def create_employee(request):
    """Crear nuevo empleado con usuario"""
    if not user_is_admin(request.user):
        messages.error(request, 'No tienes permisos para crear empleados')
        return redirect('employee_dashboard')
    
    if request.method == 'POST':
        form = EmployeeForm(request.POST, user=request.user)
        if form.is_valid():
            employee = form.save()
            messages.success(request, f'Empleado "{employee.full_name}" creado exitosamente')
            messages.info(request, f'Usuario creado: {form.cleaned_data["user_email"]}')
            messages.info(request, f'Contraseña: {form.user_password}')
            return redirect('employee_list')
    else:
        form = EmployeeForm(user=request.user)
    
    context = {
        'form': form,
        'title': 'Nuevo Empleado',
        'ciudades_por_provincia': json.dumps(CIUDADES_POR_PROVINCIA)
    }
    return render(request, 'admin/employee_form.html', context)


@login_required
def edit_employee(request, employee_id):
    """Editar empleado existente"""
    if not user_is_admin(request.user):
        messages.error(request, 'No tienes permisos para editar empleados')
        return redirect('employee_dashboard')
    
    employee = get_object_or_404(Employee, id=employee_id)
    
    # Verificar permisos
    if not (request.user.is_superuser or request.user.role == CustomUser.Role.SUPERUSER):
        if employee.company not in request.user.managed_companies.all():
            messages.error(request, 'No tienes permisos para editar este empleado')
            return redirect('employee_list')
    
    if request.method == 'POST':
        form = EmployeeForm(request.POST, instance=employee, user=request.user)
        if form.is_valid():
            employee = form.save()
            messages.success(request, f'Empleado "{employee.full_name}" actualizado exitosamente')
            if form.cleaned_data.get('user_password'):
                messages.info(request, f'Nueva contraseña: {form.user_password}')
            return redirect('employee_list')
    else:
        form = EmployeeForm(instance=employee, user=request.user)
    
    context = {
        'form': form,
        'employee': employee,
        'title': 'Editar Empleado',
        'ciudades_por_provincia': json.dumps(CIUDADES_POR_PROVINCIA)
    }
    return render(request, 'admin/employee_form.html', context)


@login_required
def company_progress(request, company_id):
    """Vista de progreso de evaluaciones para la empresa"""
    if not user_is_admin(request.user):
        messages.error(request, 'No tienes permisos para acceder a esta página')
        return redirect('employee_dashboard')
    
    company = get_object_or_404(Company, id=company_id)
    
    # Verificar permisos usando función de seguridad
    if not validate_company_access(request.user, company):
        messages.error(request, 'No tienes permisos para ver esta empresa')
        return redirect('admin_dashboard')
    
    employees = Employee.objects.filter(company=company)
    current_year = datetime.now().year
    
    # Obtener progreso por empleado
    progress_data = []
    for employee in employees:
        evaluation = Evaluation.objects.filter(
            employee=employee,
            year=current_year
        ).first()
        
        if evaluation:
            total_questions = Question.objects.count()
            answered = Response.objects.filter(evaluation=evaluation).count()
            progress_percentage = (answered / total_questions * 100) if total_questions > 0 else 0
            status = 'Completada' if evaluation.is_complete else 'En progreso'
        else:
            progress_percentage = 0
            status = 'No iniciada'
        
        progress_data.append({
            'employee': employee,
            'evaluation': evaluation,
            'progress_percentage': progress_percentage,
            'status': status
        })
    
    context = {
        'company': company,
        'progress_data': progress_data,
        'current_year': current_year,
        'total_employees': employees.count(),
        'completed': sum(1 for p in progress_data if p['status'] == 'Completada'),
        'in_progress': sum(1 for p in progress_data if p['status'] == 'En progreso'),
        'not_started': sum(1 for p in progress_data if p['status'] == 'No iniciada'),
    }
    
    return render(request, 'admin/company_progress.html', context)


def generate_password():
    """Genera una contraseña segura"""
    alphabet = string.ascii_letters + string.digits
    return ''.join(secrets.choice(alphabet) for i in range(12))


@login_required
def employee_detail(request, employee_id):
    """Vista de detalle del empleado"""
    if not user_is_admin(request.user):
        messages.error(request, 'No tienes permisos para acceder a esta página')
        return redirect('employee_dashboard')
    
    employee = get_object_or_404(Employee, id=employee_id)
    
    # Verificar permisos usando función de seguridad
    if not validate_employee_access(request.user, employee):
        messages.error(request, 'No tienes permisos para ver este empleado')
        return redirect('employee_list')
    
    # Obtener evaluaciones del empleado
    evaluations = Evaluation.objects.filter(employee=employee).order_by('-year', '-date_started')
    
    # Obtener contraseña desencriptada solo si es superusuario
    employee_password_decrypted = None
    if employee.user and employee.user.stored_password:
        employee_password_decrypted = employee.user.get_stored_password(request.user)
    
    context = {
        'employee': employee,
        'evaluations': evaluations,
        'employee_password_decrypted': employee_password_decrypted,
    }
    
    return render(request, 'admin/employee_detail.html', context)


@login_required
def admin_view_evaluation(request, evaluation_id):
    """Vista para que superusuarios puedan ver el cuestionario de cualquier evaluación"""
    if not (request.user.is_superuser or request.user.role == CustomUser.Role.SUPERUSER):
        messages.error(request, 'No tienes permisos para acceder a esta página')
        return redirect('admin_dashboard')
    
    evaluation = get_object_or_404(Evaluation, id=evaluation_id)
    employee = evaluation.employee
    
    # Obtener todas las preguntas
    questions = Question.objects.all().order_by('number')
    
    # Obtener respuestas existentes
    existing_responses = {
        r.question_id: r.answer 
        for r in Response.objects.filter(evaluation=evaluation)
    }
    
    # Preparar datos de provincias y ciudades
    from apps.presentation.utils.ecuador_data import PROVINCIAS_ECUADOR, CIUDADES_POR_PROVINCIA, get_todas_las_ciudades
    provincias = PROVINCIAS_ECUADOR
    todas_las_ciudades = get_todas_las_ciudades()
    
    # Si hay una provincia seleccionada, obtener sus ciudades
    provincia_seleccionada = evaluation.province or employee.province
    ciudades = CIUDADES_POR_PROVINCIA.get(provincia_seleccionada, todas_las_ciudades) if provincia_seleccionada else todas_las_ciudades
    
    # Preparar ciudades por provincia para JavaScript
    ciudades_por_provincia = {provincia: ciudades for provincia, ciudades in CIUDADES_POR_PROVINCIA.items()}
    
    context = {
        'evaluation': evaluation,
        'employee': employee,
        'questions': questions,
        'existing_responses': existing_responses,
        'completion_percentage': evaluation.completion_percentage,
        'provincias': provincias,
        'ciudades': ciudades,
        'ciudades_por_provincia': ciudades_por_provincia,
        'is_superuser_view': True,  # Flag para indicar que es vista de superusuario
    }
    
    return render(request, 'employee/take_evaluation.html', context)


@login_required
def admin_view_empty_evaluation(request):
    """Vista para que superusuarios puedan ver el cuestionario vacío para revisión"""
    if not (request.user.is_superuser or request.user.role == CustomUser.Role.SUPERUSER):
        messages.error(request, 'No tienes permisos para acceder a esta página')
        return redirect('admin_dashboard')
    
    # Obtener el primer empleado disponible o crear uno temporal para la vista
    employee = Employee.objects.first()
    if not employee:
        messages.error(request, 'No hay empleados en el sistema. Por favor, crea al menos un empleado primero.')
        return redirect('superuser_dashboard')
    
    current_year = datetime.now().year
    
    # Crear una evaluación temporal solo para visualización (no se guarda)
    # Usamos get_or_create para evitar duplicados si se accede múltiples veces
    evaluation, created = Evaluation.objects.get_or_create(
        employee=employee,
        year=current_year,
        status=Evaluation.Status.DRAFT,
        defaults={'year': current_year}
    )
    
    # Obtener todas las preguntas
    questions = Question.objects.all().order_by('number')
    
    # No hay respuestas existentes (cuestionario vacío)
    existing_responses = {}
    
    # Preparar datos de provincias y ciudades
    from apps.presentation.utils.ecuador_data import PROVINCIAS_ECUADOR, CIUDADES_POR_PROVINCIA, get_todas_las_ciudades
    provincias = PROVINCIAS_ECUADOR
    todas_las_ciudades = get_todas_las_ciudades()
    
    # Obtener ciudades por defecto
    ciudades = todas_las_ciudades
    
    # Preparar ciudades por provincia para JavaScript
    ciudades_por_provincia = {provincia: ciudades for provincia, ciudades in CIUDADES_POR_PROVINCIA.items()}
    
    context = {
        'evaluation': evaluation,
        'employee': employee,
        'questions': questions,
        'existing_responses': existing_responses,
        'completion_percentage': 0,  # Cuestionario vacío
        'provincias': provincias,
        'ciudades': ciudades,
        'ciudades_por_provincia': ciudades_por_provincia,
        'is_superuser_view': True,  # Flag para indicar que es vista de superusuario
        'is_empty_preview': True,  # Flag adicional para indicar que es vista vacía
    }
    
    return render(request, 'employee/take_evaluation.html', context)


@login_required
def employee_toggle_status(request, employee_id):
    """Cambiar estado del empleado (AJAX)"""
    if not user_is_admin(request.user):
        return JsonResponse({'success': False, 'error': 'No tienes permisos'}, status=403)
    
    employee = get_object_or_404(Employee, id=employee_id)
    
    # Verificar permisos
    if not (request.user.is_superuser or request.user.role == CustomUser.Role.SUPERUSER):
        if employee.company not in request.user.managed_companies.all():
            return JsonResponse({'success': False, 'error': 'No tienes permisos para editar este empleado'}, status=403)
    
    if request.method == 'POST':
        try:
            new_status = request.POST.get('status', '').lower()
            if new_status == 'active':
                employee.user.is_active = True
            elif new_status == 'inactive':
                employee.user.is_active = False
            else:
                return JsonResponse({'success': False, 'error': 'Estado inválido'}, status=400)
            
            employee.user.save()
            return JsonResponse({
                'success': True,
                'is_active': employee.user.is_active,
                'status_text': 'Activo' if employee.user.is_active else 'Inactivo'
            })
        except Exception as e:
            return JsonResponse({'success': False, 'error': str(e)}, status=500)
    
    return JsonResponse({'success': False, 'error': 'Método no permitido'}, status=405)


@login_required
@require_http_methods(["POST"])
def employee_delete(request, employee_id):
    """Eliminar empleado"""
    if not user_is_admin(request.user):
        messages.error(request, 'No tienes permisos para eliminar empleados')
        return redirect('employee_dashboard')
    
    employee = get_object_or_404(Employee, id=employee_id)
    
    # Verificar permisos
    if not (request.user.is_superuser or request.user.role == CustomUser.Role.SUPERUSER):
        if employee.company not in request.user.managed_companies.all():
            messages.error(request, 'No tienes permisos para eliminar este empleado')
            return redirect('employee_list')
    
    employee_name = employee.full_name
    company_id = employee.company.id if employee.company else None
    
    # Eliminar empleado (esto también eliminará el usuario asociado por CASCADE)
    employee.delete()
    
    messages.success(request, f'Empleado "{employee_name}" eliminado exitosamente')
    
    if company_id:
        return redirect('employee_list_by_company', company_id=company_id)
    return redirect('employee_list')


@login_required
@require_http_methods(["POST"])
def employee_bulk_delete(request):
    """Eliminación masiva de empleados"""
    if not user_is_admin(request.user):
        messages.error(request, 'No tienes permisos para eliminar empleados')
        return redirect('employee_dashboard')
    
    employee_ids = request.POST.getlist('selected_employees')
    
    if not employee_ids:
        messages.warning(request, 'No se seleccionaron empleados para eliminar')
        return redirect('employee_list')
    
    employees = Employee.objects.filter(id__in=employee_ids)
    
    # Verificar permisos
    if not (request.user.is_superuser or request.user.role == CustomUser.Role.SUPERUSER):
        companies = request.user.managed_companies.all()
        employees = employees.filter(company__in=companies)
    
    count = employees.count()
    employees.delete()
    
    messages.success(request, f'{count} empleado(s) eliminado(s) exitosamente')
    
    company_id = request.POST.get('company_id')
    if company_id:
        return redirect('employee_list_by_company', company_id=company_id)
    return redirect('employee_list')


@login_required
def employee_export(request, company_id=None):
    """Exportar empleados a Excel o CSV"""
    if not user_is_admin(request.user):
        messages.error(request, 'No tienes permisos para exportar datos')
        return redirect('employee_dashboard')
    
    export_format = request.GET.get('format', 'excel')
    company = None
    
    # Obtener empleados
    if company_id:
        company = get_object_or_404(Company, id=company_id)
        employees = Employee.objects.filter(company=company)
    else:
        if request.user.is_superuser or request.user.role == CustomUser.Role.SUPERUSER:
            employees = Employee.objects.all()
        else:
            companies = request.user.managed_companies.all()
            employees = Employee.objects.filter(company__in=companies)
    
    employees = employees.select_related('user', 'company').order_by('last_name', 'first_name')
    
    if export_format == 'csv':
        return export_employees_csv(employees, company)
    else:
        return export_employees_excel(employees, company)


def export_employees_excel(employees, company=None):
    """Exportar empleados a Excel"""
    workbook = openpyxl.Workbook()
    sheet = workbook.active
    sheet.title = "Empleados"
    
    # Encabezados
    headers = [
        'Cédula', 'Nombres', 'Apellidos', 'Email', 'Empresa', 'Cargo', 'Área Interna', 'Área ERP',
        'Fecha Nacimiento', 'Edad', 'Género', 'Etnia', 'Educación',
        'Fecha Ingreso', 'Años Antigüedad', 'Provincia', 'Ciudad', 'Estado', 'Fecha Creación'
    ]
    
    # Estilo para encabezados
    header_fill = PatternFill(start_color="366092", end_color="366092", fill_type="solid")
    header_font = Font(bold=True, color="FFFFFF")
    
    for col_num, header in enumerate(headers, 1):
        cell = sheet.cell(row=1, column=col_num)
        cell.value = header
        cell.fill = header_fill
        cell.font = header_font
        cell.alignment = Alignment(horizontal="center", vertical="center")
    
    # Datos
    for row_num, employee in enumerate(employees, 2):
        sheet.cell(row=row_num, column=1).value = employee.identification
        sheet.cell(row=row_num, column=2).value = employee.first_name
        sheet.cell(row=row_num, column=3).value = employee.last_name
        sheet.cell(row=row_num, column=4).value = employee.user.email
        sheet.cell(row=row_num, column=5).value = employee.company.name
        sheet.cell(row=row_num, column=6).value = employee.position
        sheet.cell(row=row_num, column=7).value = employee.area
        sheet.cell(row=row_num, column=8).value = employee.get_work_area_erp_display()
        sheet.cell(row=row_num, column=9).value = employee.date_of_birth.strftime('%d/%m/%Y')
        sheet.cell(row=row_num, column=10).value = employee.age
        sheet.cell(row=row_num, column=11).value = employee.get_gender_display()
        sheet.cell(row=row_num, column=12).value = employee.get_ethnicity_display()
        sheet.cell(row=row_num, column=13).value = employee.get_education_level_display()
        sheet.cell(row=row_num, column=14).value = employee.hire_date.strftime('%d/%m/%Y')
        sheet.cell(row=row_num, column=15).value = employee.years_of_experience
        sheet.cell(row=row_num, column=16).value = employee.province
        sheet.cell(row=row_num, column=17).value = employee.city
        sheet.cell(row=row_num, column=18).value = 'Activo' if employee.user.is_active else 'Inactivo'
        sheet.cell(row=row_num, column=19).value = employee.created_at.strftime('%d/%m/%Y %H:%M')
    
    # Ajustar ancho de columnas
    for col_num in range(1, len(headers) + 1):
        column_letter = get_column_letter(col_num)
        sheet.column_dimensions[column_letter].width = 20
    
    # Crear respuesta
    response = HttpResponse(
        content_type='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
    )
    filename = f"empleados_{company.name if company else 'todos'}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.xlsx"
    response['Content-Disposition'] = f'attachment; filename="{filename}"'
    
    workbook.save(response)
    return response


def export_employees_csv(employees, company=None):
    """Exportar empleados a CSV"""
    response = HttpResponse(content_type='text/csv; charset=utf-8')
    filename = f"empleados_{company.name if company else 'todos'}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.csv"
    response['Content-Disposition'] = f'attachment; filename="{filename}"'
    
    writer = csv.writer(response)
    writer.writerow([
        'Cédula', 'Nombres', 'Apellidos', 'Email', 'Empresa', 'Cargo', 'Área Interna', 'Área ERP',
        'Fecha Nacimiento', 'Edad', 'Género', 'Etnia', 'Educación',
        'Fecha Ingreso', 'Años Antigüedad', 'Provincia', 'Ciudad', 'Estado', 'Fecha Creación'
    ])
    
    for employee in employees:
        writer.writerow([
            employee.identification,
            employee.first_name,
            employee.last_name,
            employee.user.email,
            employee.company.name,
            employee.position,
            employee.area,
            employee.get_work_area_erp_display(),
            employee.date_of_birth.strftime('%d/%m/%Y'),
            employee.age,
            employee.get_gender_display(),
            employee.get_ethnicity_display(),
            employee.get_education_level_display(),
            employee.hire_date.strftime('%d/%m/%Y'),
            employee.years_of_experience,
            employee.province,
            employee.city,
            'Activo' if employee.user.is_active else 'Inactivo',
            employee.created_at.strftime('%d/%m/%Y %H:%M')
        ])
    
    return response


def import_companies_from_excel(sheet):
    """Importa empresas desde Excel"""
    count = 0
    for row in sheet.iter_rows(min_row=2, values_only=True):
        if not row[0]:  # Si no hay nombre, saltar
            continue
        
        try:
            admin_email = row[7] if len(row) > 7 and row[7] else f"admin.{row[1]}@empresa.com"
            password = generate_password()
            
            admin_user, created = CustomUser.objects.get_or_create(
                email=admin_email,
                defaults={
                    'role': CustomUser.Role.COMPANY_ADMIN,
                    'is_staff': True,
                    'is_active': True
                }
            )
            if created:
                admin_user.set_password(password)
                admin_user.set_stored_password(password)  # Encriptar y almacenar contraseña
                admin_user.save()
            elif not admin_user.stored_password:
                # Usuario existente sin contraseña almacenada: generar una nueva y guardarla
                admin_user.set_password(password)
                admin_user.set_stored_password(password)  # Encriptar y almacenar contraseña
                admin_user.save()
            
            Company.objects.get_or_create(
                ruc=str(row[1]),
                defaults={
                    'name': row[0],
                    'address': row[2] or '',
                    'phone': str(row[3]) if row[3] else '',
                    'email': row[4] or admin_email,
                    'city': row[5] or '',
                    'province': row[6] or '',
                    'admin': admin_user
                }
            )
            count += 1
        except Exception as e:
            print(f"Error importando empresa: {e}")
            continue
    
    return count


def import_employees_from_excel(sheet):
    """Importa empleados desde Excel"""
    count = 0
    for row in sheet.iter_rows(min_row=2, values_only=True):
        if not row[0]:  # Si no hay nombre, saltar
            continue
        
        try:
            # Buscar empresa por RUC o nombre
            company_ruc = str(row[10]) if len(row) > 10 and row[10] else None
            company_name = row[11] if len(row) > 11 and row[11] else None
            
            company = None
            if company_ruc:
                company = Company.objects.filter(ruc=company_ruc).first()
            elif company_name:
                company = Company.objects.filter(name=company_name).first()
            
            if not company:
                continue
            
            user_email = row[12] if len(row) > 12 and row[12] else f"{str(row[0]).lower().replace(' ', '.')}@empresa.com"
            
            # Leer contraseña del Excel si está presente (columna 14, índice 13)
            password_from_excel = None
            if len(row) > 13 and row[13]:
                password_str = str(row[13]).strip()
                if password_str:  # Solo si hay contenido después de quitar espacios
                    password_from_excel = password_str
            
            employee_user, created = CustomUser.objects.get_or_create(
                email=user_email,
                defaults={
                    'role': CustomUser.Role.EMPLOYEE,
                    'is_active': True
                }
            )
            
            # Manejar contraseña:
            # - Si es usuario nuevo: siempre asignar contraseña (del Excel o generada)
            # - Si es usuario existente y hay contraseña en Excel: actualizar contraseña
            # - Si es usuario existente y NO hay contraseña en Excel: generar nueva y guardarla en stored_password
            if created:
                # Usuario nuevo: usar contraseña del Excel o generar una
                password = password_from_excel if password_from_excel else generate_password()
                employee_user.set_password(password)
                employee_user.set_stored_password(password)  # Encriptar y almacenar contraseña
                employee_user.save()
            elif password_from_excel:
                # Usuario existente pero se proporciona contraseña desde Excel: actualizar
                employee_user.set_password(password_from_excel)
                employee_user.set_stored_password(password_from_excel)  # Encriptar y almacenar contraseña
                employee_user.save()
            elif not employee_user.stored_password:
                # Usuario existente sin contraseña almacenada: generar una nueva y guardarla
                # Esto es útil para usuarios creados antes de agregar el campo stored_password
                password = generate_password()
                employee_user.set_password(password)
                employee_user.set_stored_password(password)  # Encriptar y almacenar contraseña
                employee_user.save()
            # Si el usuario ya existe y tiene stored_password, mantener la contraseña actual
            
            # Parsear fecha de nacimiento
            date_of_birth = None
            if row[3]:
                if isinstance(row[3], datetime):
                    date_of_birth = row[3].date()
                else:
                    try:
                        date_of_birth = datetime.strptime(str(row[3]), '%Y-%m-%d').date()
                    except:
                        try:
                            date_of_birth = datetime.strptime(str(row[3]), '%d/%m/%Y').date()
                        except:
                            date_of_birth = datetime(1990, 1, 1).date()
            else:
                date_of_birth = datetime(1990, 1, 1).date()
            
            # Calcular fecha de ingreso desde años de experiencia (si está en Excel)
            from datetime import date
            from dateutil.relativedelta import relativedelta
            
            hire_date_value = date(2020, 1, 1)  # Default
            if row[9]:
                if isinstance(row[9], date):
                    hire_date_value = row[9]
                else:
                    try:
                        years_exp = int(row[9])
                        hire_date_value = date.today() - relativedelta(years=years_exp)
                    except (ValueError, TypeError):
                        pass
            
            # Leer work_area_erp (columna 7, índice 7 - nuevo campo)
            work_area_erp_value = 'ADMINISTRATIVA'  # Default
            if len(row) > 7 and row[7]:
                work_area_erp_str = str(row[7]).upper().strip()
                # Validar que sea una opción válida
                valid_areas = ['ADMINISTRATIVA', 'OPERATIVA', 'COMERCIAL', 'LOGISTICA', 'PRODUCCION', 'TRANSFORMACION', 'SOPORTE', 'TECNOLOGIA_INNOVACION']
                if work_area_erp_str in valid_areas:
                    work_area_erp_value = work_area_erp_str
            
            Employee.objects.get_or_create(
                identification=str(row[2]),
                defaults={
                    'user': employee_user,
                    'company': company,
                    'first_name': row[0],
                    'last_name': row[1],
                    'date_of_birth': date_of_birth,
                    'gender': row[4] or 'M',
                    'ethnicity': row[5] or 'MESTIZO',
                    'area': row[6] or '',
                    'work_area_erp': work_area_erp_value,
                    'position': row[8] or '',  # Ahora position está en columna 8
                    'education_level': row[9] or 'UNIVERSITARIO',  # education_level en columna 9
                    'hire_date': hire_date_value,
                    'province': company.province,
                    'city': company.city
                }
            )
            count += 1
        except Exception as e:
            print(f"Error importando empleado: {e}")
            continue
    
    return count


def import_evaluations_from_excel(sheet):
    """Importa evaluaciones históricas desde Excel"""
    count = 0
    for row in sheet.iter_rows(min_row=2, values_only=True):
        if not row[0]:  # Si no hay identificación, saltar
            continue
        
        try:
            employee = Employee.objects.filter(identification=str(row[0])).first()
            if not employee:
                continue
            
            year = int(row[1]) if row[1] else datetime.now().year
            
            evaluation, created = Evaluation.objects.get_or_create(
                employee=employee,
                year=year,
                defaults={
                    'status': Evaluation.Status.COMPLETED,
                    'date_completed': row[2] if isinstance(row[2], datetime) else datetime.now(),
                }
            )
            
            # Importar respuestas (asumiendo que están en columnas 3-60)
            if created:
                questions = Question.objects.all().order_by('number')
                for idx, question in enumerate(questions):
                    answer_value = row[3 + idx] if len(row) > 3 + idx and row[3 + idx] else None
                    if answer_value:
                        try:
                            answer_int = int(answer_value)
                            if 1 <= answer_int <= 4:
                                Response.objects.create(
                                    evaluation=evaluation,
                                    question=question,
                                    answer=answer_int
                                )
                        except:
                            continue
                
                # Calcular riesgos solo si hay respuestas
                if Response.objects.filter(evaluation=evaluation).count() > 0:
                    calculator = RiskCalculatorService()
                    calculator.calculate_evaluation_risk(evaluation)
            
            count += 1
        except Exception as e:
            print(f"Error importando evaluación: {e}")
            continue
    
    return count


@login_required
def download_import_template(request):
    """Descargar plantilla Excel para importación"""
    if not user_is_admin(request.user):
        messages.error(request, 'No tienes permisos para descargar plantillas')
        return redirect('admin_dashboard')
    
    template_type = request.GET.get('type', 'employees')
    
    workbook = openpyxl.Workbook()
    sheet = workbook.active
    sheet.title = "Plantilla"
    
    if template_type == 'employees':
        # Encabezados para empleados
        headers = [
            'Nombres', 'Apellidos', 'Cédula', 'Fecha Nacimiento (YYYY-MM-DD)', 
            'Género (M/F)', 'Etnia (MESTIZO/INDIGENA/AFROECUATORIANO/BLANCO/MONTUBIO/OTRO)',
            'Área Interna', 'Área ERP (ADMINISTRATIVA/OPERATIVA/COMERCIAL/LOGISTICA/PRODUCCION/TRANSFORMACION/SOPORTE/TECNOLOGIA_INNOVACION)', 
            'Cargo', 'Educación (PRIMARIA/SECUNDARIA/TECNICO/UNIVERSITARIO/POSTGRADO)',
            'Fecha Ingreso (YYYY-MM-DD) o Años Antigüedad', 'RUC Empresa', 'Nombre Empresa', 'Email Empleado', 'Contraseña (Opcional)'
        ]
        
        # Estilo para encabezados
        header_fill = PatternFill(start_color="366092", end_color="366092", fill_type="solid")
        header_font = Font(bold=True, color="FFFFFF", size=11)
        
        for col_num, header in enumerate(headers, 1):
            cell = sheet.cell(row=1, column=col_num)
            cell.value = header
            cell.fill = header_fill
            cell.font = header_font
            cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
        
        # Ejemplo de datos
        example_row = [
            'Juan', 'Pérez', '1234567890', '1990-01-15',
            'M', 'MESTIZO',
            'Ventas y Marketing', 'COMERCIAL', 'Analista', 'UNIVERSITARIO',
            '2018-03-15', '1234567890001', 'Empresa Demo S.A.', 'juan.perez@empresa.com', ''
        ]
        
        # Nota explicativa en la siguiente fila
        note_row = [
            '* Columna 7 "Área ERP": Use las opciones MDT (ADMINISTRATIVA, OPERATIVA, COMERCIAL, LOGISTICA, PRODUCCION, TRANSFORMACION, SOPORTE, TECNOLOGIA_INNOVACION)', '', '', '', '', '', '', '', '',
            'Columna 10 puede ser FECHA (2018-03-15) o NÚMERO de años (5)', '', '', '', ''
        ]
        
        for col_num, value in enumerate(example_row, 1):
            cell = sheet.cell(row=2, column=col_num)
            cell.value = value
            cell.fill = PatternFill(start_color="E7F3FF", end_color="E7F3FF", fill_type="solid")
            cell.font = Font(italic=True, color="666666")
        
        # Ajustar ancho de columnas (incluye nueva columna Área ERP)
        column_widths = [15, 15, 15, 25, 15, 30, 20, 35, 20, 30, 25, 15, 25, 30, 25]
        for col_num, width in enumerate(column_widths, 1):
            column_letter = get_column_letter(col_num)
            sheet.column_dimensions[column_letter].width = width
        
        filename = "plantilla_importacion_empleados.xlsx"
        
    elif template_type == 'companies':
        # Encabezados para empresas
        headers = [
            'Nombre', 'RUC', 'Dirección', 'Teléfono', 'Email', 
            'Ciudad', 'Provincia', 'Email Administrador'
        ]
        
        header_fill = PatternFill(start_color="366092", end_color="366092", fill_type="solid")
        header_font = Font(bold=True, color="FFFFFF", size=11)
        
        for col_num, header in enumerate(headers, 1):
            cell = sheet.cell(row=1, column=col_num)
            cell.value = header
            cell.fill = header_fill
            cell.font = header_font
            cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
        
        example_row = [
            'Empresa Ejemplo S.A.', '1234567890001', 'Av. Principal 123', '0999999999',
            'contacto@empresa.com', 'Quito', 'Pichincha', 'admin@empresa.com'
        ]
        
        for col_num, value in enumerate(example_row, 1):
            cell = sheet.cell(row=2, column=col_num)
            cell.value = value
            cell.fill = PatternFill(start_color="E7F3FF", end_color="E7F3FF", fill_type="solid")
            cell.font = Font(italic=True, color="666666")
        
        column_widths = [25, 15, 30, 15, 25, 15, 15, 25]
        for col_num, width in enumerate(column_widths, 1):
            column_letter = get_column_letter(col_num)
            sheet.column_dimensions[column_letter].width = width
        
        filename = "plantilla_importacion_empresas.xlsx"
    else:
        # Plantilla para evaluaciones
        headers = ['Cédula Empleado', 'Año', 'Fecha Completada (YYYY-MM-DD)']
        # Agregar columnas para las 58 preguntas
        for i in range(1, 59):
            headers.append(f'Pregunta {i} (1-4)')
        
        header_fill = PatternFill(start_color="366092", end_color="366092", fill_type="solid")
        header_font = Font(bold=True, color="FFFFFF", size=10)
        
        for col_num, header in enumerate(headers, 1):
            cell = sheet.cell(row=1, column=col_num)
            cell.value = header
            cell.fill = header_fill
            cell.font = header_font
            cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
        
        filename = "plantilla_importacion_evaluaciones.xlsx"
    
    response = HttpResponse(
        content_type='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
    )
    response['Content-Disposition'] = f'attachment; filename="{filename}"'
    
    workbook.save(response)
    return response


@login_required
def bulk_import(request):
    """Importación masiva desde Excel"""
    if not (request.user.is_superuser or request.user.role == CustomUser.Role.SUPERUSER):
        messages.error(request, 'Solo los superusuarios pueden importar datos masivos')
        return redirect('admin_dashboard')
    
    if request.method == 'POST':
        form = BulkImportForm(request.POST, request.FILES)
        if form.is_valid():
            excel_file = request.FILES['excel_file']
            import_type = form.cleaned_data['import_type']
            
            try:
                workbook = openpyxl.load_workbook(excel_file)
                sheet = workbook.active
                
                with transaction.atomic():
                    if import_type == 'companies':
                        count = import_companies_from_excel(sheet)
                        messages.success(request, f'{count} empresas importadas exitosamente')
                    elif import_type == 'employees':
                        count = import_employees_from_excel(sheet)
                        messages.success(request, f'{count} empleados importados exitosamente')
                    elif import_type == 'evaluations':
                        count = import_evaluations_from_excel(sheet)
                        messages.success(request, f'{count} evaluaciones importadas exitosamente')
                
                return redirect('admin_dashboard')
            except Exception as e:
                messages.error(request, f'Error al importar: {str(e)}')
    else:
        form = BulkImportForm()
    
    return render(request, 'admin/bulk_import.html', {'form': form})


@login_required
def download_anonymous_evaluations_excel(request, company_id=None):
    """Descarga Excel con estructura MDT: datos demográficos + nivel de riesgo por dimensión"""
    if not user_is_admin(request.user):
        messages.error(request, 'No tienes permisos para acceder a esta página')
        return redirect('employee_dashboard')
    
    # Filtrar por empresa
    company = None
    if company_id:
        company = get_object_or_404(Company, id=company_id)
        evaluations = Evaluation.objects.filter(
            employee__company=company,
            status=Evaluation.Status.COMPLETED
        ).select_related('employee')
    else:
        if request.user.is_superuser or request.user.role == CustomUser.Role.SUPERUSER:
            evaluations = Evaluation.objects.filter(status=Evaluation.Status.COMPLETED).select_related('employee')
        else:
            companies = request.user.managed_companies.all()
            if companies.count() == 1:
                company = companies.first()
            evaluations = Evaluation.objects.filter(
                employee__company__in=companies,
                status=Evaluation.Status.COMPLETED
            ).select_related('employee')
    
    # Obtener año del filtro
    # Validar año de forma segura
    try:
        year_str = request.GET.get('year', str(datetime.now().year))
        year = validate_year(year_str)
        year_str = str(year)
    except (ValueError, TypeError, ValidationError):
        year = datetime.now().year
        year_str = str(year)
    try:
        year = int(year_str)
    except (ValueError, TypeError):
        year = datetime.now().year
    evaluations = evaluations.filter(year=year)
    
    # Obtener todas las dimensiones ordenadas
    from apps.infrastructure.models import Dimension
    all_dimensions = Dimension.objects.all().order_by('order')
    
    # Obtener todas las preguntas ordenadas
    all_questions = Question.objects.all().order_by('id')
    
    # Crear el workbook - ESTRUCTURA MDT COMPLETA
    workbook = openpyxl.Workbook()
    ws = workbook.active
    ws.title = "Tabulación Completa"
    
    # ==================== ENCABEZADOS COMPLETOS ====================
    # Fila 1: Encabezados principales
    headers_row1 = [
        'N° Cuestionario',  # Col A - NUEVO
        'Fecha Evaluación',  # Col B - NUEVO
        'Ciudad',  # Col C
        'Nivel más alto de instrucción',  # Col D
        'Antigüedad, años de experiencia dentro de la empresa',  # Col E
        'Edad del trabajador',  # Col F
        'Sexo del trabajador',  # Col G
        'Departamento laboral',  # Col H
    ]
    
    # Agregar encabezados de dimensiones
    for dim in all_dimensions:
        headers_row1.append(f"DIMENSIÓN {dim.order}. {dim.name.upper()}")
    
    # Agregar encabezados de preguntas (P1, P2, ... P58)
    for i in range(1, len(all_questions) + 1):
        headers_row1.append(f"P{i}")
    
    # Estilos
    id_fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")  # Azul oscuro para ID y Fecha
    id_font = Font(bold=True, size=10, color="FFFFFF")
    header_fill = PatternFill(start_color="D9E1F2", end_color="D9E1F2", fill_type="solid")  # Azul claro para demográficos
    header_font = Font(bold=True, size=10)
    dim_fill = PatternFill(start_color="FFC000", end_color="FFC000", fill_type="solid")  # Naranja para dimensiones
    dim_font = Font(bold=True, size=9, color="000000")
    question_fill = PatternFill(start_color="70AD47", end_color="70AD47", fill_type="solid")  # Verde para preguntas
    question_font = Font(bold=True, size=8, color="FFFFFF")
    
    # Aplicar encabezados
    num_demographic_cols = 8  # N° Cuestionario, Fecha, Ciudad, Nivel, Antigüedad, Edad, Sexo, Departamento
    num_dimension_cols = len(all_dimensions)  # 14 dimensiones
    
    for col_num, header in enumerate(headers_row1, 1):
        cell = ws.cell(row=1, column=col_num)
        cell.value = header
        
        if col_num <= 2:  # N° Cuestionario y Fecha
            cell.fill = id_fill
            cell.font = id_font
        elif col_num <= num_demographic_cols:  # Columnas demográficas
            cell.fill = header_fill
            cell.font = header_font
        elif col_num <= num_demographic_cols + num_dimension_cols:  # Columnas de dimensiones
            cell.fill = dim_fill
            cell.font = dim_font
        else:  # Columnas de preguntas (P1, P2, ...)
            cell.fill = question_fill
            cell.font = question_font
        
        cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
    
    # ==================== DATOS ====================
    # Colores para niveles de riesgo (como en la imagen MDT)
    color_bajo = PatternFill(start_color="C6E0B4", end_color="C6E0B4", fill_type="solid")  # Verde claro
    color_medio = PatternFill(start_color="FFE699", end_color="FFE699", fill_type="solid")  # Amarillo claro
    color_alto = PatternFill(start_color="F4B084", end_color="F4B084", fill_type="solid")  # Naranja claro
    
    for row_num, evaluation in enumerate(evaluations, 2):
        # Columna A: N° Cuestionario
        ws.cell(row=row_num, column=1).value = f"EVAL-{str(evaluation.id).zfill(4)}"
        ws.cell(row=row_num, column=1).alignment = Alignment(horizontal="center", vertical="center")
        
        # Columna B: Fecha de evaluación
        fecha_eval = evaluation.date_completed if evaluation.date_completed else evaluation.date_started
        ws.cell(row=row_num, column=2).value = fecha_eval.strftime('%d/%m/%Y') if fecha_eval else "N/A"
        ws.cell(row=row_num, column=2).alignment = Alignment(horizontal="center", vertical="center")
        
        # Columna C: Ciudad
        ws.cell(row=row_num, column=3).value = evaluation.city or "N/A"
        
        # Columna D: Nivel de instrucción
        ws.cell(row=row_num, column=4).value = evaluation.get_education_level_display()
        
        # Columna E: Antigüedad
        ws.cell(row=row_num, column=5).value = evaluation.get_experience_range_display()
        
        # Columna F: Edad
        ws.cell(row=row_num, column=6).value = evaluation.get_age_range_display()
        
        # Columna G: Sexo
        ws.cell(row=row_num, column=7).value = evaluation.get_gender_display()
        
        # Columna H: Departamento laboral
        ws.cell(row=row_num, column=8).value = evaluation.get_work_area_display()
        
        # Obtener resultados de riesgo para esta evaluación
        risk_results = RiskResult.objects.filter(evaluation=evaluation).select_related('dimension')
        risk_dict = {result.dimension.id: result for result in risk_results}
        
        # Columnas I en adelante: Una columna por cada dimensión (14 dimensiones)
        for col_offset, dim in enumerate(all_dimensions):
            col_num = 9 + col_offset  # Empieza en columna 9 (I)
            cell = ws.cell(row=row_num, column=col_num)
            
            if dim.id in risk_dict:
                risk_result = risk_dict[dim.id]
                cell.value = f"RIESGO {risk_result.risk_level}"
                
                # Aplicar color según nivel de riesgo
                if risk_result.risk_level == 'BAJO':
                    cell.fill = color_bajo
                elif risk_result.risk_level == 'MEDIO':
                    cell.fill = color_medio
                elif risk_result.risk_level == 'ALTO':
                    cell.fill = color_alto
                
                cell.font = Font(bold=False, size=9)
                cell.alignment = Alignment(horizontal="center", vertical="center")
            else:
                cell.value = "N/A"
        
        # Después de las dimensiones: agregar las 58 respuestas (P1, P2, ... P58)
        responses = Response.objects.filter(evaluation=evaluation).order_by('question__id')
        response_dict = {resp.question.id: resp.answer for resp in responses}
        
        # Columnas después de las dimensiones (columna 23 en adelante si hay 14 dimensiones)
        start_col_questions = 9 + len(all_dimensions)
        for col_offset, question in enumerate(all_questions):
            col_num = start_col_questions + col_offset
            answer = response_dict.get(question.id, '')
            cell = ws.cell(row=row_num, column=col_num)
            cell.value = answer
            cell.alignment = Alignment(horizontal="center", vertical="center")
    
    # ==================== AJUSTAR ANCHOS DE COLUMNAS ====================
    ws.column_dimensions['A'].width = 18  # N° Cuestionario
    ws.column_dimensions['B'].width = 15  # Fecha
    ws.column_dimensions['C'].width = 15  # Ciudad
    ws.column_dimensions['D'].width = 25  # Nivel instrucción
    ws.column_dimensions['E'].width = 20  # Antigüedad
    ws.column_dimensions['F'].width = 15  # Edad
    ws.column_dimensions['G'].width = 12  # Sexo
    ws.column_dimensions['H'].width = 25  # Departamento
    
    # Dimensiones (columnas I-V, 14 dimensiones)
    for i in range(9, 9 + len(all_dimensions)):
        ws.column_dimensions[get_column_letter(i)].width = 20
    
    # Preguntas (columnas después de dimensiones, 58 preguntas)
    start_col_questions = 9 + len(all_dimensions)
    for i in range(start_col_questions, start_col_questions + len(all_questions)):
        ws.column_dimensions[get_column_letter(i)].width = 8
    
    # ==================== AGREGAR HOJA DE RESUMEN ====================
    ws_resumen = workbook.create_sheet(title="Resumen")
    
    # Título
    ws_resumen.cell(row=1, column=1).value = "RESUMEN DE EVALUACIONES DE RIESGO PSICOSOCIAL"
    ws_resumen.cell(row=1, column=1).font = Font(bold=True, size=14, color="1F4E78")
    ws_resumen.merge_cells('A1:D1')
    ws_resumen.cell(row=1, column=1).alignment = Alignment(horizontal="center", vertical="center")
    
    # Información general
    company_name = company.name if company else "Todas las empresas"
    
    info_data = [
        ['Empresa:', company_name],
        ['Año:', year],
        ['Total Evaluaciones:', evaluations.count()],
        ['Fecha:', datetime.now().strftime('%d/%m/%Y %H:%M')],
    ]
    
    row = 3
    for label, value in info_data:
        ws_resumen.cell(row=row, column=1).value = label
        ws_resumen.cell(row=row, column=1).font = Font(bold=True)
        ws_resumen.cell(row=row, column=2).value = value
        row += 1
    
    # Calcular distribución por nivel de riesgo
    total = evaluations.count()
    
    row += 1
    ws_resumen.cell(row=row, column=1).value = "DISTRIBUCIÓN POR NIVEL DE RIESGO"
    ws_resumen.cell(row=row, column=1).font = Font(bold=True, size=12)
    ws_resumen.merge_cells(f'A{row}:C{row}')
    
    row += 1
    headers = ['Nivel', 'Cantidad', 'Porcentaje']
    for col, header in enumerate(headers, 1):
        cell = ws_resumen.cell(row=row, column=col)
        cell.value = header
        cell.font = Font(bold=True)
        cell.fill = PatternFill(start_color="D9E1F2", end_color="D9E1F2", fill_type="solid")
        cell.alignment = Alignment(horizontal="center", vertical="center")
    
    # Calcular estadísticas por dimensión
    bajo_total = 0
    medio_total = 0
    alto_total = 0
    
    for evaluation in evaluations:
        risk_results = RiskResult.objects.filter(evaluation=evaluation)
        for result in risk_results:
            if result.risk_level == 'BAJO':
                bajo_total += 1
            elif result.risk_level == 'MEDIO':
                medio_total += 1
            elif result.risk_level == 'ALTO':
                alto_total += 1
    
    grand_total = bajo_total + medio_total + alto_total
    
    # Fila BAJO
    row += 1
    ws_resumen.cell(row=row, column=1).value = "RIESGO BAJO"
    ws_resumen.cell(row=row, column=1).fill = color_bajo
    ws_resumen.cell(row=row, column=2).value = bajo_total
    ws_resumen.cell(row=row, column=3).value = f"{(bajo_total/grand_total*100):.1f}%" if grand_total > 0 else "0%"
    ws_resumen.cell(row=row, column=2).alignment = Alignment(horizontal="center")
    ws_resumen.cell(row=row, column=3).alignment = Alignment(horizontal="center")
    
    # Fila MEDIO
    row += 1
    ws_resumen.cell(row=row, column=1).value = "RIESGO MEDIO"
    ws_resumen.cell(row=row, column=1).fill = color_medio
    ws_resumen.cell(row=row, column=2).value = medio_total
    ws_resumen.cell(row=row, column=3).value = f"{(medio_total/grand_total*100):.1f}%" if grand_total > 0 else "0%"
    ws_resumen.cell(row=row, column=2).alignment = Alignment(horizontal="center")
    ws_resumen.cell(row=row, column=3).alignment = Alignment(horizontal="center")
    
    # Fila ALTO
    row += 1
    ws_resumen.cell(row=row, column=1).value = "RIESGO ALTO"
    ws_resumen.cell(row=row, column=1).fill = color_alto
    ws_resumen.cell(row=row, column=2).value = alto_total
    ws_resumen.cell(row=row, column=3).value = f"{(alto_total/grand_total*100):.1f}%" if grand_total > 0 else "0%"
    ws_resumen.cell(row=row, column=2).alignment = Alignment(horizontal="center")
    ws_resumen.cell(row=row, column=3).alignment = Alignment(horizontal="center")
    
    # Nota de confidencialidad
    row += 3
    ws_resumen.cell(row=row, column=1).value = "NOTA DE CONFIDENCIALIDAD:"
    ws_resumen.cell(row=row, column=1).font = Font(bold=True, color="FF0000", size=11)
    ws_resumen.merge_cells(f'A{row}:D{row}')
    
    row += 1
    ws_resumen.cell(row=row, column=1).value = "Los datos en este archivo mantienen el ANONIMATO de los participantes. No se incluyen nombres, cédulas, emails ni información que permita la identificación directa de las personas evaluadas."
    ws_resumen.merge_cells(f'A{row}:D{row}')
    ws_resumen.cell(row=row, column=1).alignment = Alignment(wrap_text=True)
    ws_resumen.row_dimensions[row].height = 30
    
    # Ajustar anchos
    ws_resumen.column_dimensions['A'].width = 25
    ws_resumen.column_dimensions['B'].width = 20
    ws_resumen.column_dimensions['C'].width = 15
    ws_resumen.column_dimensions['D'].width = 50
    
    # Preparar respuesta HTTP
    response = HttpResponse(
        content_type='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
    )
    company_name_safe = company.name if company else "Todas_las_empresas"
    filename = f"Tabulacion_Riesgo_Psicosocial_{company_name_safe.replace(' ', '_')}_{year}.xlsx"
    response['Content-Disposition'] = f'attachment; filename="{filename}"'
    
    workbook.save(response)
    return response


@login_required
def preview_pptx_report(request, company_id=None):
    """Vista HTML optimizada para previsualización y captura de screenshots para PowerPoint"""
    if not user_is_admin(request.user):
        messages.error(request, 'No tienes permisos para acceder a esta página')
        return redirect('employee_dashboard')
    
    # Reutilizar TODA la lógica de evaluation_results
    # Filtrar por empresa
    company = None
    if company_id:
        company = get_object_or_404(Company, id=company_id)
        evaluations = Evaluation.objects.filter(
            employee__company=company,
            status=Evaluation.Status.COMPLETED
        ).select_related('employee')
    else:
        if request.user.is_superuser or request.user.role == CustomUser.Role.SUPERUSER:
            evaluations = Evaluation.objects.filter(status=Evaluation.Status.COMPLETED).select_related('employee')
        else:
            companies = request.user.managed_companies.all()
            if companies.count() == 1:
                company = companies.first()
            evaluations = Evaluation.objects.filter(
                employee__company__in=companies,
                status=Evaluation.Status.COMPLETED
            ).select_related('employee')
    
    # Obtener año del filtro
    # Validar año de forma segura
    try:
        year_str = request.GET.get('year', str(datetime.now().year))
        year = validate_year(year_str)
        year_str = str(year)
    except (ValueError, TypeError, ValidationError):
        year = datetime.now().year
        year_str = str(year)
    try:
        year = int(year_str)
    except (ValueError, TypeError):
        year = datetime.now().year
    evaluations = evaluations.filter(year=year)
    
    # Calcular todas las estadísticas (igual que evaluation_results)
    risk_results = RiskResult.objects.filter(evaluation__in=evaluations).select_related('dimension')
    all_dimensions = Dimension.objects.all().order_by('order')
    total_evaluations_count = evaluations.count()
    
    # Calcular dimension_stats
    dimension_stats = {}
    dimension_stats_list = []
    
    for dim in all_dimensions:
        bajo_count = risk_results.filter(dimension=dim, risk_level='BAJO').count()
        medio_count = risk_results.filter(dimension=dim, risk_level='MEDIO').count()
        alto_count = risk_results.filter(dimension=dim, risk_level='ALTO').count()
        
        bajo_percentage = round((bajo_count / total_evaluations_count * 100), 1) if total_evaluations_count > 0 else 0
        medio_percentage = round((medio_count / total_evaluations_count * 100), 1) if total_evaluations_count > 0 else 0
        alto_percentage = round((alto_count / total_evaluations_count * 100), 1) if total_evaluations_count > 0 else 0
        
        dimension_stats[dim.name] = {
            'bajo_count': bajo_count,
            'medio_count': medio_count,
            'alto_count': alto_count,
            'bajo_percentage': bajo_percentage,
            'medio_percentage': medio_percentage,
            'alto_percentage': alto_percentage,
        }
        
        dimension_stats_list.append({
            'dimension_name': dim.name,
            'bajo_count': bajo_count,
            'medio_count': medio_count,
            'alto_count': alto_count,
            'bajo_percentage': bajo_percentage,
            'medio_percentage': medio_percentage,
            'alto_percentage': alto_percentage,
        })
    
    # Calcular riesgo global (por puntuación total)
    from apps.application.services.risk_calculator import RiskCalculatorService
    calculator = RiskCalculatorService()
    
    global_score_counts = {'bajo': 0, 'medio': 0, 'alto': 0}
    for evaluation in evaluations:
        global_score = calculator.calculate_global_score(evaluation)
        global_risk_level = calculator.calculate_global_risk_level(global_score)
        if global_risk_level == 'BAJO':
            global_score_counts['bajo'] += 1
        elif global_risk_level == 'MEDIO':
            global_score_counts['medio'] += 1
        elif global_risk_level == 'ALTO':
            global_score_counts['alto'] += 1
    
    global_score_percentages = {
        'bajo': round((global_score_counts['bajo'] / total_evaluations_count * 100), 1) if total_evaluations_count > 0 else 0,
        'medio': round((global_score_counts['medio'] / total_evaluations_count * 100), 1) if total_evaluations_count > 0 else 0,
        'alto': round((global_score_counts['alto'] / total_evaluations_count * 100), 1) if total_evaluations_count > 0 else 0,
    }
    
    # Determinar riesgo general
    if global_score_counts['alto'] > global_score_counts['medio'] and global_score_counts['alto'] > global_score_counts['bajo']:
        overall_risk = 'ALTO'
    elif global_score_counts['medio'] >= global_score_counts['bajo'] and global_score_counts['medio'] >= global_score_counts['alto']:
        overall_risk = 'MEDIO'
    else:
        overall_risk = 'BAJO'
    
    # Recomendaciones
    from apps.application.services.employer_recommendations_service import EmployerRecommendationsService
    employer_recs_service = EmployerRecommendationsService()
    general_recommendations = employer_recs_service.get_general_recommendations_by_level(overall_risk)
    
    # Calcular datos demográficos (MISMA LÓGICA que evaluation_results)
    def calculate_demographic_stats_preview(field_name, display_method=None):
        """Calcula estadísticas demográficas - IDÉNTICA a evaluation_results"""
        # Obtener valores únicos con su conteo
        field_values = evaluations.values(field_name).annotate(count=Count('id')).order_by(field_name)
        
        stats = []
        for item in field_values:
            value = item[field_name]
            if not value:  # Saltar valores vacíos
                continue
            
            # Filtrar evaluaciones por este valor
            filtered_evals = evaluations.filter(**{field_name: value})
            count = filtered_evals.count()
            
            # Obtener label legible
            if display_method:
                # Obtener una evaluación de ejemplo para usar get_FOO_display()
                sample_eval = filtered_evals.first()
                if sample_eval:
                    label = getattr(sample_eval, display_method)()
                else:
                    label = value
            else:
                label = value
            
            # Calcular riesgo por dimensión para este grupo
            group_risk_results = RiskResult.objects.filter(
                evaluation__in=filtered_evals
            ).select_related('dimension')
            
            # Inicializar contadores por dimensión
            dimension_risks = {}
            for dim in all_dimensions:
                dimension_risks[dim.name] = {
                    'BAJO': 0,
                    'MEDIO': 0,
                    'ALTO': 0,
                    'porcentaje_bajo': 0.0,
                    'porcentaje_medio': 0.0,
                    'porcentaje_alto': 0.0,
                }
            
            # Contar resultados
            for result in group_risk_results:
                dim_name = result.dimension.name
                if dim_name in dimension_risks:
                    dimension_risks[dim_name][result.risk_level] += 1
            
            # Calcular porcentajes
            if count > 0:
                for dim_name, risks in dimension_risks.items():
                    risks['porcentaje_bajo'] = round((risks['BAJO'] / count) * 100, 1)
                    risks['porcentaje_medio'] = round((risks['MEDIO'] / count) * 100, 1)
                    risks['porcentaje_alto'] = round((risks['ALTO'] / count) * 100, 1)
            
            stats.append({
                'value': value,
                'label': label,
                'count': count,
                'dimension_risks': dimension_risks
            })
        
        return stats
    
    demographic_data = {
        'gender': calculate_demographic_stats_preview('gender', 'get_gender_display'),
        'age_range': calculate_demographic_stats_preview('age_range', 'get_age_range_display'),
        'work_area': calculate_demographic_stats_preview('work_area', 'get_work_area_display'),
        'education_level': calculate_demographic_stats_preview('education_level', 'get_education_level_display'),
        'experience_range': calculate_demographic_stats_preview('experience_range', 'get_experience_range_display'),
        'ethnicity': calculate_demographic_stats_preview('ethnicity', 'get_ethnicity_display'),
        'province': calculate_demographic_stats_preview('province'),
        'city': calculate_demographic_stats_preview('city'),
    }
    
    context = {
        'company': company,
        'year': year,
        'total_evaluations_count': total_evaluations_count,
        'all_dimensions': all_dimensions,
        'dimension_stats': dimension_stats,
        'dimension_stats_list': dimension_stats_list,
        'global_score_counts': global_score_counts,
        'global_score_percentages': global_score_percentages,
        'overall_risk': overall_risk,
        'general_recommendations': general_recommendations,
        'demographic_data': demographic_data,
    }
    
    return render(request, 'admin/pptx_preview.html', context)


@login_required
def pptx_progress_page(request, company_id=None):
    """Muestra página de progreso mientras se genera el PowerPoint"""
    if not user_is_admin(request.user):
        messages.error(request, 'No tienes permisos para acceder a esta página')
        return redirect('employee_dashboard')
    
    year = request.GET.get('year', str(datetime.now().year))
    
    if company_id:
        company = get_object_or_404(Company, id=company_id)
        company_name = company.name
        generate_url = reverse('generate_pptx_with_screenshots_by_company', kwargs={'company_id': company_id}) + f'?year={year}'
        back_url = reverse('evaluation_results')
    else:
        company_name = "Resultados_Globales"
        generate_url = reverse('generate_pptx_with_screenshots') + f'?year={year}'
        back_url = reverse('evaluation_results')
    
    filename = f"Informe_Riesgo_Psicosocial_{company_name.replace(' ', '_')}_{year}.pptx"
    
    context = {
        'company_name': company_name,
        'year': year,
        'generate_url': generate_url,
        'filename': filename,
        'back_url': back_url
    }
    
    return render(request, 'admin/pptx_progress.html', context)


@login_required
def pdf_preview_page(request, company_id=None):
    """Página de preview para generar PDF profesional con Playwright"""
    if not user_is_admin(request.user):
        messages.error(request, 'No tienes permisos para acceder a esta página')
        return redirect('employee_dashboard')
    
    # Reutilizar la lógica de evaluation_results para obtener todos los datos
    year = request.GET.get('year', str(datetime.now().year))
    
    if company_id:
        company = get_object_or_404(Company, id=company_id)
        if not request.user.is_superuser and request.user.role == CustomUser.Role.COMPANY_ADMIN:
            # Verificar si el usuario tiene employee antes de acceder
            if hasattr(request.user, 'employee') and request.user.employee:
                if request.user.employee.company_id != company.id:
                    messages.error(request, 'No tienes permisos para acceder a esta empresa')
                    return redirect('admin_dashboard')
        evaluations = Evaluation.objects.filter(employee__company=company, year=year, status='COMPLETED')
    else:
        company = None
        evaluations = Evaluation.objects.filter(year=year, status='COMPLETED')
    
    total_evaluations_count = evaluations.count()
    
    if total_evaluations_count == 0:
        messages.warning(request, f'No hay evaluaciones completadas para el año {year}')
        return redirect('evaluation_results')
    
    # Obtener todas las dimensiones
    all_dimensions = Dimension.objects.all().order_by('order')
    
    # Calcular estadísticas por dimensión
    dimension_stats = {}
    for dimension in all_dimensions:
        results = RiskResult.objects.filter(
            evaluation__in=evaluations,
            dimension=dimension
        )
        
        total = results.count()
        if total > 0:
            bajo = results.filter(risk_level='BAJO').count()
            medio = results.filter(risk_level='MEDIO').count()
            alto = results.filter(risk_level='ALTO').count()
            
            dimension_stats[dimension.name] = {
                'bajo_count': bajo,
                'medio_count': medio,
                'alto_count': alto,
                'bajo_percentage': round((bajo / total) * 100, 1),
                'medio_percentage': round((medio / total) * 100, 1),
                'alto_percentage': round((alto / total) * 100, 1),
                'predominant': 'BAJO' if bajo >= medio and bajo >= alto else ('MEDIO' if medio >= alto else 'ALTO')
            }
    
    # Calcular puntaje global y niveles de riesgo
    from apps.application.services.risk_calculator import RiskCalculatorService
    
    risk_calculator = RiskCalculatorService()
    global_risk_evaluations = {'BAJO': 0, 'MEDIO': 0, 'ALTO': 0}
    for evaluation in evaluations:
        global_score = risk_calculator.calculate_global_score(evaluation)
        global_risk_level = risk_calculator.calculate_global_risk_level(global_score)
        global_risk_evaluations[global_risk_level] += 1
    
    global_score_counts = {
        'bajo': global_risk_evaluations['BAJO'],
        'medio': global_risk_evaluations['MEDIO'],
        'alto': global_risk_evaluations['ALTO']
    }
    
    global_score_percentages = {
        'bajo': round((global_score_counts['bajo'] / total_evaluations_count) * 100, 1) if total_evaluations_count > 0 else 0,
        'medio': round((global_score_counts['medio'] / total_evaluations_count) * 100, 1) if total_evaluations_count > 0 else 0,
        'alto': round((global_score_counts['alto'] / total_evaluations_count) * 100, 1) if total_evaluations_count > 0 else 0
    }
    
    # Determinar riesgo predominante global
    max_count = max(global_score_counts.values())
    if global_score_counts['bajo'] == max_count:
        overall_risk = 'BAJO'
    elif global_score_counts['medio'] == max_count:
        overall_risk = 'MEDIO'
    else:
        overall_risk = 'ALTO'
    
    # ==================== RECOMENDACIONES ====================
    from apps.application.services.employer_recommendations_service import EmployerRecommendationsService
    employer_service = EmployerRecommendationsService()
    
    # 1. Recomendaciones Estratégicas Generales (basadas en el riesgo global)
    strategic_recommendations = employer_service.get_general_recommendations_by_level(overall_risk)
    
    # 2. Recomendaciones por Dimensión (Plan de Acción detallado)
    dimension_recommendations = {}
    for dim_name, stats in dimension_stats.items():
        predominant_level = stats['predominant']
        
        # Agregar recomendaciones para TODOS los niveles (BAJO, MEDIO, ALTO)
        recommendations = employer_service.get_recommendations_by_dimension(dim_name, predominant_level)
        if recommendations:
            dimension_recommendations[dim_name] = {
                'level': predominant_level,
                'count': stats[f'{predominant_level.lower()}_count'],
                'percentage': stats[f'{predominant_level.lower()}_percentage'],
                'recommendations': recommendations
            }
    
    # ==================== CALCULAR DATOS DEMOGRÁFICOS COMPLETOS ====================
    # Reutilizar la función de evaluation_results para obtener TODOS los datos demográficos
    def calculate_demographic_stats_pdf(field_name, display_method=None):
        """Calcula estadísticas por campo demográfico incluyendo recomendaciones - SIN DUPLICACIONES"""
        stats = []
        seen_values = set()  # Para evitar duplicaciones
        
        # Obtener valores únicos del campo en las evaluaciones - CONVERTIR A LISTA Y USAR SET
        unique_values = list(evaluations.values_list(field_name, flat=True).distinct())
        
        # Eliminar duplicados y None
        unique_values = list(set([v for v in unique_values if v]))
        
        for value in unique_values:
            if not value or value in seen_values:
                continue
            
            seen_values.add(value)  # Marcar como visto
            
            # Filtrar evaluaciones por este valor
            filtered_evals = evaluations.filter(**{field_name: value})
            count = filtered_evals.count()
            
            if count == 0:
                continue
            
            # Obtener label
            if display_method:
                sample_eval = filtered_evals.first()
                label = getattr(sample_eval, display_method)() if hasattr(sample_eval, display_method) else value
            else:
                label = value
            
            # Verificar que no hayamos agregado este label antes
            if any(s['label'] == label for s in stats):
                continue
            
            # Calcular riesgo por dimensión para este grupo
            group_risk_results = RiskResult.objects.filter(
                evaluation__in=filtered_evals
            ).select_related('dimension')
            
            # Inicializar contadores por dimensión
            dimension_risks = {}
            for dim in all_dimensions:
                dimension_risks[dim.name] = {
                    'BAJO': 0,
                    'MEDIO': 0,
                    'ALTO': 0,
                    'porcentaje_bajo': 0.0,
                    'porcentaje_medio': 0.0,
                    'porcentaje_alto': 0.0,
                }
            
            # Contar resultados
            for result in group_risk_results:
                dim_name = result.dimension.name
                if dim_name in dimension_risks:
                    dimension_risks[dim_name][result.risk_level] += 1
            
            # Calcular porcentajes
            if count > 0:
                for dim_name, risks in dimension_risks.items():
                    risks['porcentaje_bajo'] = round((risks['BAJO'] / count) * 100, 1)
                    risks['porcentaje_medio'] = round((risks['MEDIO'] / count) * 100, 1)
                    risks['porcentaje_alto'] = round((risks['ALTO'] / count) * 100, 1)
            
            # Obtener recomendaciones para este grupo
            group_recommendations = {}
            for dim_name, risks in dimension_risks.items():
                # Determinar nivel predominante
                if risks['ALTO'] > 0:
                    level = 'ALTO'
                    rec_count = risks['ALTO']
                    percentage = risks['porcentaje_alto']
                elif risks['MEDIO'] > 0:
                    level = 'MEDIO'
                    rec_count = risks['MEDIO']
                    percentage = risks['porcentaje_medio']
                else:
                    continue  # No agregar recomendaciones para bajo
                
                # Obtener recomendaciones del servicio
                recs = employer_service.get_recommendations_by_dimension(dim_name, level)
                if recs:
                    group_recommendations[dim_name] = {
                        'level': level,
                        'count': rec_count,
                        'percentage': percentage,
                        'recommendations': recs
                    }
            
            stats.append({
                'value': value,
                'label': label,
                'count': count,
                'dimension_risks': dimension_risks,
                'recommendations': group_recommendations
            })
        
        return stats
    
    # Calcular todos los análisis demográficos
    demographic_data = {
        'gender': calculate_demographic_stats_pdf('gender', 'get_gender_display'),
        'age_range': calculate_demographic_stats_pdf('age_range', 'get_age_range_display'),
        'work_area': calculate_demographic_stats_pdf('work_area', 'get_work_area_display'),
        'education_level': calculate_demographic_stats_pdf('education_level', 'get_education_level_display'),
        'experience_range': calculate_demographic_stats_pdf('experience_range', 'get_experience_range_display'),
        'ethnicity': calculate_demographic_stats_pdf('ethnicity', 'get_ethnicity_display'),
        'province': calculate_demographic_stats_pdf('province'),
        'city': calculate_demographic_stats_pdf('city'),
    }
    
    # ==================== CALCULAR CUMPLIMIENTO ====================
    def calculate_completion_stats_pdf(employee_field, evaluation_field=None, display_method=None):
        """Calcula estadísticas de cumplimiento - SIN DUPLICACIONES"""
        from apps.infrastructure.models.company import Employee
        
        if evaluation_field is None:
            evaluation_field = employee_field
        
        completion_stats = []
        seen_labels = set()  # Para evitar duplicaciones
        
        # Para age_range y experience_range, usar las opciones del modelo
        if employee_field == 'age_range':
            unique_values = [choice[0] for choice in Employee.AgeRange.choices]
        elif employee_field == 'experience_range':
            unique_values = [choice[0] for choice in Employee.ExperienceRange.choices]
        else:
            # Obtener valores únicos de employees - CONVERTIR A LISTA Y USAR SET
            if company:
                unique_values = list(Employee.objects.filter(company=company).values_list(employee_field, flat=True).distinct())
            else:
                unique_values = list(Employee.objects.values_list(employee_field, flat=True).distinct())
        
        # Eliminar duplicados y None
        unique_values = list(set([v for v in unique_values if v]))
        
        for value in unique_values:
            if not value:
                continue
            
            # Contar programados
            if employee_field in ['age_range', 'experience_range']:
                # Para propiedades calculadas, iterar manualmente
                if company:
                    programmed_employees = [emp for emp in Employee.objects.filter(company=company) 
                                          if getattr(emp, employee_field) == value]
                else:
                    programmed_employees = [emp for emp in Employee.objects.all() 
                                          if getattr(emp, employee_field) == value]
                programmed = len(programmed_employees)
            else:
                if company:
                    programmed = Employee.objects.filter(company=company, **{employee_field: value}).count()
                else:
                    programmed = Employee.objects.filter(**{employee_field: value}).count()
            
            # Contar completados
            completed = evaluations.filter(**{evaluation_field: value}).count()
            
            # Calcular porcentaje
            percentage = round((completed / programmed * 100), 1) if programmed > 0 else 0
            
            # Obtener label
            if display_method:
                if employee_field in ['age_range', 'experience_range'] and 'programmed_employees' in locals() and programmed_employees:
                    label = getattr(programmed_employees[0], display_method)() if hasattr(programmed_employees[0], display_method) else value
                else:
                    from apps.infrastructure.models.evaluation import Evaluation
                    sample_eval = evaluations.filter(**{evaluation_field: value}).first()
                    if sample_eval and hasattr(sample_eval, display_method):
                        label = getattr(sample_eval, display_method)()
                    else:
                        label = value
            else:
                label = value
            
            # EVITAR DUPLICACIONES: Solo agregar si el label no se ha visto antes
            if label not in seen_labels:
                seen_labels.add(label)
                completion_stats.append({
                    'label': label,
                    'programmed': programmed,
                    'completed': completed,
                    'percentage': percentage
                })
        
        return completion_stats
    
    # Calcular datos de cumplimiento
    compliance_data = {
        'gender': calculate_completion_stats_pdf('gender', 'gender', 'get_gender_display'),
        'age_range': calculate_completion_stats_pdf('age_range', 'age_range', 'get_age_range_display'),
        'work_area': calculate_completion_stats_pdf('area', 'work_area', 'get_work_area_display'),
        'education': calculate_completion_stats_pdf('education_level', 'education_level', 'get_education_level_display'),
        'experience': calculate_completion_stats_pdf('experience_range', 'experience_range', 'get_experience_range_display'),
        'ethnicity': calculate_completion_stats_pdf('ethnicity', 'ethnicity', 'get_ethnicity_display'),
        'province': calculate_completion_stats_pdf('province', 'province'),
        'city': calculate_completion_stats_pdf('city', 'city'),
    }
    
    context = {
        'company': company,
        'year': year,
        'total_evaluations_count': total_evaluations_count,
        'all_dimensions': all_dimensions,
        'dimension_stats': dimension_stats,
        'global_score_counts': global_score_counts,
        'global_score_percentages': global_score_percentages,
        'overall_risk': overall_risk,
        'strategic_recommendations': strategic_recommendations,  # Recomendaciones estratégicas generales
        'dimension_recommendations': dimension_recommendations,  # Recomendaciones por dimensión
        'demographic_data': demographic_data,
        'compliance_data': compliance_data,
    }
    
    return render(request, 'admin/pdf_preview.html', context)


@login_required
def pdf_progress_page(request, company_id=None):
    """Muestra página de progreso mientras se genera el PDF"""
    if not user_is_admin(request.user):
        messages.error(request, 'No tienes permisos para acceder a esta página')
        return redirect('employee_dashboard')
    
    year = request.GET.get('year', str(datetime.now().year))
    
    if company_id:
        company = get_object_or_404(Company, id=company_id)
        company_name = company.name
        generate_url = reverse('generate_pdf_playwright_by_company', kwargs={'company_id': company_id}) + f'?year={year}'
        back_url = reverse('evaluation_results')
    else:
        company_name = "Resultados_Globales"
        generate_url = reverse('generate_pdf_playwright') + f'?year={year}'
        back_url = reverse('evaluation_results')
    
    filename = f"Informe_Riesgo_Psicosocial_{company_name.replace(' ', '_')}_{year}.pdf"
    
    context = {
        'company_name': company_name,
        'year': year,
        'generate_url': generate_url,
        'filename': filename,
        'back_url': back_url
    }
    
    return render(request, 'admin/pdf_progress.html', context)


@login_required
def generate_pdf_with_playwright(request, company_id=None):
    """Genera PDF profesional capturando la página de preview con Playwright"""
    if not user_is_admin(request.user):
        messages.error(request, 'No tienes permisos para acceder a esta página')
        return redirect('employee_dashboard')
    
    try:
        from playwright.sync_api import sync_playwright
    except ImportError:
        messages.error(request, 'Playwright no está instalado. Ejecuta: pip install playwright && playwright install chromium')
        return redirect('evaluation_results')
    
    year = request.GET.get('year', str(datetime.now().year))
    
    if company_id:
        company = get_object_or_404(Company, id=company_id)
        company_name = company.name
    else:
        company_name = "Resultados_Globales"
    
    # Construir URL de preview
    preview_url = request.build_absolute_uri(
        reverse('pdf_preview_by_company' if company_id else 'pdf_preview', 
                kwargs={'company_id': company_id} if company_id else None) + f'?year={year}'
    )
    
    print(f"🌐 Navegando a preview PDF: {preview_url}")
    
    # Generar PDF con Playwright
    with sync_playwright() as p:
        browser = p.chromium.launch(headless=True)
        context = browser.new_context(viewport={'width': 1920, 'height': 1080})
        
        # Agregar cookie de sesión para autenticación
        session_cookie = {
            'name': 'sessionid',
            'value': request.session.session_key,
            'domain': '127.0.0.1',
            'path': '/',
            'httpOnly': True,
            'secure': False
        }
        context.add_cookies([session_cookie])
        print(f"🔐 Cookie de sesión agregada")
        
        page = context.new_page()
        
        # Navegar y esperar carga completa
        page.goto(preview_url, wait_until='load', timeout=120000)
        print("✅ Página HTML cargada")
        
        # Esperar a que los gráficos se rendericen
        print("⏳ Esperando renderizado de gráficos...")
        page.wait_for_timeout(15000)
        
        print("📄 Generando PDF...")
        # Generar PDF directamente con Playwright
        pdf_bytes = page.pdf(
            format='A4',
            print_background=True,
            margin={
                'top': '1cm',
                'right': '1cm',
                'bottom': '1cm',
                'left': '1cm'
            },
            prefer_css_page_size=True
        )
        
        browser.close()
        print(f"✅ PDF generado: {len(pdf_bytes)} bytes")
    
    # Retornar PDF
    response = HttpResponse(pdf_bytes, content_type='application/pdf')
    filename = f"Informe_Riesgo_Psicosocial_{company_name.replace(' ', '_')}_{year}.pdf"
    response['Content-Disposition'] = f'attachment; filename="{filename}"'
    
    return response


@login_required
def generate_pptx_with_screenshots(request, company_id=None):
    """Genera PowerPoint capturando screenshots de los gráficos con Playwright"""
    if not user_is_admin(request.user):
        messages.error(request, 'No tienes permisos para acceder a esta página')
        return redirect('employee_dashboard')
    
    try:
        from playwright.sync_api import sync_playwright
    except ImportError:
        messages.error(request, 'Playwright no está instalado. Ejecuta: pip install playwright && playwright install chromium')
        return redirect('evaluation_results')
    
    # Construir URL de preview
    if company_id:
        preview_url = request.build_absolute_uri(
            reverse('preview_pptx_report_by_company', kwargs={'company_id': company_id})
        )
        company = get_object_or_404(Company, id=company_id)
        company_name = company.name
    else:
        preview_url = request.build_absolute_uri(reverse('preview_pptx_report'))
        company_name = "Resultados_Globales"
    
    # Agregar año a la URL
    year = request.GET.get('year', str(datetime.now().year))
    preview_url += f'?year={year}'
    
    # Capturar screenshots con Playwright
    screenshots = {}
    
    with sync_playwright() as p:
        browser = p.chromium.launch(headless=True)
        context = browser.new_context(viewport={'width': 1920, 'height': 1080})
        
        # ⭐ IMPORTANTE: Pasar cookies de sesión para autenticación
        # Obtener el sessionid del usuario actual
        session_cookie = {
            'name': 'sessionid',
            'value': request.session.session_key,
            'domain': '127.0.0.1',
            'path': '/',
            'httpOnly': True,
            'secure': False
        }
        context.add_cookies([session_cookie])
        print(f"🔐 Cookie de sesión agregada: {request.session.session_key[:10]}...")
        
        page = context.new_page()
        
        print(f"🌐 Navegando a: {preview_url}")
        
        # Estrategia simplificada: solo esperar 'load' y dar tiempo fijo
        page.goto(preview_url, wait_until='load', timeout=120000)
        
        print("✅ Página HTML cargada")
        
        # Esperar tiempo fijo para que Chart.js y todos los gráficos se carguen
        # Esto es más confiable que esperar por elementos específicos
        print("⏳ Esperando 15 segundos para que se carguen y rendericen todos los gráficos...")
        page.wait_for_timeout(15000)  # 15 segundos fijos
        
        print("✅ Tiempo de espera completado, iniciando captura de screenshots...")
        
        # Debugging: verificar qué hay en la página
        try:
            # Primero verificar si hay ALGÚN elemento en la página
            body_content = page.content()
            print(f"📄 Contenido HTML recibido: {len(body_content)} caracteres")
            
            # Intentar contar slides de varias formas
            slide_count = page.locator('.slide').count()
            print(f"📊 Slides encontradas con '.slide': {slide_count}")
            
            if slide_count == 0:
                print("⚠️ No se encontraron elementos con clase 'slide'")
                # Intentar buscar por ID
                all_divs = page.locator('div[id^="slide-"]').count()
                print(f"📊 Elementos con id 'slide-*': {all_divs}")
                
                # Guardar HTML para debugging
                with open('debug_page.html', 'w', encoding='utf-8') as f:
                    f.write(body_content)
                print("💾 HTML guardado en debug_page.html para inspección")
            
            # Capturar cada slide por índice
            for idx in range(slide_count):
                try:
                    print(f"🎯 Intentando capturar slide {idx+1}/{slide_count}...")
                    element = page.locator('.slide').nth(idx)
                    
                    # Verificar que el elemento existe
                    is_visible = element.is_visible()
                    print(f"   Visible: {is_visible}")
                    
                    screenshot_bytes = element.screenshot()
                    screenshots[f'slide-{idx+1}'] = screenshot_bytes
                    print(f"✅ Capturada slide {idx+1}/{slide_count} ({len(screenshot_bytes)} bytes)")
                except Exception as e:
                    print(f"❌ Error capturando slide {idx+1}: {e}")
                    import traceback
                    traceback.print_exc()
                    continue
            
            print(f"🎉 Total capturadas: {len(screenshots)} slides")
        except Exception as e:
            print(f"❌ Error general capturando slides: {e}")
            import traceback
            traceback.print_exc()
        
        browser.close()
    
    # VALIDAR que se capturaron slides
    if not screenshots:
        print("❌ ERROR: No se capturaron screenshots. El PowerPoint estaría vacío.")
        messages.error(request, 'No se pudieron capturar las gráficas. Intenta nuevamente.')
        return redirect('evaluation_results')
    
    print(f"📦 Creando PowerPoint con {len(screenshots)} slides...")
    
    # Crear PowerPoint con las imágenes
    from pptx import Presentation
    from pptx.util import Inches
    from io import BytesIO
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # Ordenar slides por número
    sorted_screenshots = sorted(screenshots.items(), key=lambda x: int(x[0].split('-')[1]))
    
    for idx, (slide_id, img_bytes) in enumerate(sorted_screenshots, 1):
        try:
            slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Agregar imagen
            img_stream = BytesIO(img_bytes)
            left = Inches(0)
            top = Inches(0)
            width = Inches(10)
            slide.shapes.add_picture(img_stream, left, top, width=width)
            print(f"✅ Slide {idx} agregada al PowerPoint")
        except Exception as e:
            print(f"❌ Error agregando slide {idx} al PowerPoint: {e}")
    
    # Guardar y retornar
    pptx_stream = BytesIO()
    prs.save(pptx_stream)
    pptx_stream.seek(0)
    
    response = HttpResponse(
        pptx_stream,
        content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation'
    )
    filename = f"Informe_Riesgo_Psicosocial_{company_name.replace(' ', '_')}_{year}.pptx"
    response['Content-Disposition'] = f'attachment; filename="{filename}"'
    
    return response
